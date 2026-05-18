"""コイケ酸商向け pptx デザインシステム.

調査結果 (tmp/pptx_best_practices_research.md) に基づくカラー/フォント/サイズ/
レイアウト 5 パターンの定数 + ヘルパー関数。

Usage:
    from _pptx_design_system import (
        DS, slide_title_bar, add_text, add_rect, add_bullets,
        slide_footer, video_placeholder, set_speaker_notes,
    )
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class DS:
    """Design System 定数。"""
    # サイズ (16:9 フル HD)
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)
    MARGIN = Inches(0.8)
    CONTENT_W = Inches(13.333) - Inches(1.6)  # 11.733 inch

    # カラー (Codex 神回 32 + portal #F97316)
    ORANGE = RGBColor(0xF9, 0x73, 0x16)        # アクセント・数字強調
    ORANGE_DARK = RGBColor(0xC2, 0x57, 0x0C)   # 見出し・補足
    ORANGE_TINT = RGBColor(0xFF, 0xF7, 0xED)   # 背景カード
    NAVY = RGBColor(0x1F, 0x29, 0x37)          # タイトル・見出し
    GRAY_DARK = RGBColor(0x37, 0x41, 0x51)     # 本文
    GRAY_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)    # 枠線
    GRAY_BG = RGBColor(0xF9, 0xFA, 0xFB)       # 補助背景
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    VIDEO_BG = RGBColor(0x1F, 0x29, 0x37)      # 動画プレースホルダ背景

    # フォント
    FONT_JP = "Yu Gothic"  # 游ゴシック (Windows/Mac/PowerPoint 全対応)

    # サイズ (Pt)
    SIZE_TITLE_BIG = 40    # 表紙・クロージング
    SIZE_TITLE = 24        # 通常スライドタイトル
    SIZE_HEADING = 18      # 見出し
    SIZE_BODY = 13         # 本文
    SIZE_SMALL = 11        # 補足
    SIZE_TINY = 9          # ページ番号・フッター
    SIZE_ROI_NUM = 42      # ROI 数字 (大字)


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """色付き矩形 (枠線オプション)。影は無効化。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text,
             size=DS.SIZE_BODY, color=DS.GRAY_DARK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.25,
             font=DS.FONT_JP):
    """テキストボックス追加 (改行は \\n で対応)。"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, bullets,
                size=DS.SIZE_BODY, color=DS.GRAY_DARK, bold=False,
                line_spacing=1.35, space_after=4):
    """箇条書きテキストボックス。"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = DS.FONT_JP
    return tb


def slide_title_bar(slide, title, subtitle=None):
    """通常スライドの上部タイトル帯 (オレンジ縦ライン + タイトル + サブタイトル)。"""
    # オレンジ縦ライン (16:9 全高)
    add_rect(slide, Inches(0), Inches(0), Inches(0.4), DS.SLIDE_H, DS.ORANGE)
    # タイトル
    add_text(slide, Inches(0.7), Inches(0.3), DS.SLIDE_W - Inches(1.4), Inches(0.55),
             title, size=DS.SIZE_TITLE, color=DS.NAVY, bold=True)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(0.85), DS.SLIDE_W - Inches(1.4), Inches(0.35),
                 subtitle, size=DS.SIZE_SMALL, color=DS.GRAY_DARK)


def slide_footer(slide, page_no, total):
    """ページ下部 (左: ブランド、右: ページ番号)。"""
    add_text(slide, Inches(0.7), DS.SLIDE_H - Inches(0.4),
             Inches(8), Inches(0.3),
             "マツモト産業㈱ 京葉営業所 / TWF2026 みどころポータル: twf2026-portal.pages.dev",
             size=DS.SIZE_TINY, color=DS.GRAY_DARK)
    add_text(slide, DS.SLIDE_W - Inches(1.5), DS.SLIDE_H - Inches(0.4),
             Inches(1.2), Inches(0.3),
             f"{page_no:02d} / {total:02d}",
             size=DS.SIZE_TINY, color=DS.ORANGE_DARK, bold=True, align=PP_ALIGN.RIGHT)


def video_placeholder(slide, left, top, width, height, label):
    """動画埋込予定枠 (16:9 推奨)。"""
    add_rect(slide, left, top, width, height, DS.VIDEO_BG, line_color=DS.ORANGE)
    add_text(slide, left, top + height/2 - Inches(0.5),
             width, Inches(0.5),
             "▶  動画埋込予定", size=22, color=DS.ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + height/2 + Inches(0.1),
             width, Inches(0.4),
             label, size=11, color=DS.WHITE, align=PP_ALIGN.CENTER)


def set_speaker_notes(slide, notes_text):
    """スピーカーノート (発表時の台本) を設定。200-400 字目安。"""
    notes = slide.notes_slide.notes_text_frame
    notes.text = notes_text


def add_layout_title_slide(prs, eyebrow, title_main, subtitle, recipient, presenter, date):
    """Layout 1: 表紙."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 上半分オレンジ
    add_rect(slide, 0, 0, DS.SLIDE_W, Inches(4.6), DS.ORANGE)
    add_rect(slide, 0, Inches(4.6), DS.SLIDE_W, DS.SLIDE_H - Inches(4.6), DS.WHITE)
    add_text(slide, Inches(0.8), Inches(0.8), DS.SLIDE_W - Inches(1.6), Inches(0.5),
             eyebrow, size=DS.SIZE_HEADING, color=DS.WHITE, bold=True)
    add_text(slide, Inches(0.8), Inches(1.5), DS.SLIDE_W - Inches(1.6), Inches(2.2),
             title_main, size=DS.SIZE_TITLE_BIG, color=DS.WHITE, bold=True, line_spacing=1.15)
    add_text(slide, Inches(0.8), Inches(3.8), DS.SLIDE_W - Inches(1.6), Inches(0.5),
             subtitle, size=DS.SIZE_HEADING, color=DS.WHITE)
    # 下半分白
    add_text(slide, Inches(0.8), Inches(5.0), DS.SLIDE_W - Inches(1.6), Inches(0.4),
             "ご提案先", size=DS.SIZE_SMALL, color=DS.GRAY_DARK)
    add_text(slide, Inches(0.8), Inches(5.35), DS.SLIDE_W - Inches(1.6), Inches(0.5),
             recipient, size=22, color=DS.NAVY, bold=True)
    add_text(slide, Inches(0.8), Inches(6.05), DS.SLIDE_W - Inches(1.6), Inches(0.4),
             "ご提案", size=DS.SIZE_SMALL, color=DS.GRAY_DARK)
    add_text(slide, Inches(0.8), Inches(6.4), DS.SLIDE_W - Inches(1.6), Inches(0.5),
             f"{presenter}  /  {date}", size=14, color=DS.NAVY)
    return slide


def add_layout_closing_slide(prs, headline, sub_message, contact_lines):
    """Layout 5: クロージング (全面オレンジ)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, DS.SLIDE_W, DS.SLIDE_H, DS.ORANGE)
    add_text(slide, Inches(0.8), Inches(1.0), DS.SLIDE_W - Inches(1.6), Inches(0.6),
             "Closing", size=14, color=DS.WHITE)
    add_text(slide, Inches(0.8), Inches(1.6), DS.SLIDE_W - Inches(1.6), Inches(2.5),
             headline, size=DS.SIZE_TITLE_BIG, color=DS.WHITE, bold=True, line_spacing=1.2)
    add_rect(slide, Inches(0.8), Inches(4.2), DS.SLIDE_W - Inches(1.6), Inches(0.05), DS.WHITE)
    add_text(slide, Inches(0.8), Inches(4.4), DS.SLIDE_W - Inches(1.6), Inches(0.5),
             sub_message, size=DS.SIZE_HEADING - 2, color=DS.WHITE)
    add_text(slide, Inches(0.8), Inches(5.3), DS.SLIDE_W - Inches(1.6), Inches(0.4),
             "ご連絡先", size=DS.SIZE_SMALL, color=DS.WHITE)
    top = Inches(5.7)
    for i, line in enumerate(contact_lines):
        bold = i == 0
        size = DS.SIZE_HEADING if i == 0 else DS.SIZE_BODY
        add_text(slide, Inches(0.8), top + Inches(0.45 * i),
                 DS.SLIDE_W - Inches(1.6), Inches(0.4),
                 line, size=size, color=DS.WHITE, bold=bold)
    return slide


def add_layout_section_divider(prs, section_no, section_title, section_desc):
    """Layout 2: 章扉 (オレンジ縦ライン + 大見出し)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Inches(0.4), DS.SLIDE_H, DS.ORANGE)
    add_text(slide, Inches(1.0), Inches(2.3), DS.SLIDE_W - Inches(1.4), Inches(0.5),
             f"Section {section_no}", size=DS.SIZE_HEADING, color=DS.ORANGE_DARK, bold=True)
    add_text(slide, Inches(1.0), Inches(2.9), DS.SLIDE_W - Inches(1.4), Inches(1.5),
             section_title, size=DS.SIZE_TITLE_BIG, color=DS.NAVY, bold=True, line_spacing=1.15)
    add_text(slide, Inches(1.0), Inches(4.8), DS.SLIDE_W - Inches(1.4), Inches(0.8),
             section_desc, size=DS.SIZE_HEADING - 2, color=DS.GRAY_DARK, line_spacing=1.4)
    return slide
