"""コイケ酸商 v2 pptx (McKinsey 風 + portal オレンジ + 数字大字).

公式 pptx skill SKILL.md 原則 + McKinsey ブランドガイド 適用.
- Sharp corners (border radius 0)
- Strong contrast (Navy text on white、white text on navy)
- Dominance: White 60% + Navy 30% + Orange 10% (sparingly)
- Dark/light sandwich: 表紙 + クロージング dark、コンテンツ light
- Large stat callouts: 60-72pt 数字
- アクセント下線禁止
- Text-only NG: 各スライド visual element
"""
import sys, io
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent))
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from _pptx_design_system_v2 import (
    DS2, add_rect, add_text, add_bullets,
    title_header, footer, set_notes,
    stat_callout, two_col_card, video_frame,
)

ROOT = Path("D:/repos/twf2026-portal")
OUT = ROOT / "tmp/コイケ酸商_TWF2026提案_v2.pptx"
TOTAL = 24


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


prs = Presentation()
prs.slide_width = DS2.SLIDE_W
prs.slide_height = DS2.SLIDE_H


# ============ Slide 01: 表紙 (Dark McKinsey 風) ============
s = new_slide(prs)
# 全面 navy (dark sandwich の上面)
add_rect(s, 0, 0, DS2.SLIDE_W, DS2.SLIDE_H, DS2.NAVY_DARK)
# 左端オレンジ縦バー (太め、表紙だけ太く)
add_rect(s, 0, 0, Inches(0.6), DS2.SLIDE_H, DS2.ORANGE)
# 右端 logo マーク代わりに矩形 (オレンジアクセント)
add_rect(s, DS2.SLIDE_W - Inches(2.2), Inches(0.6), Inches(1.6), Inches(0.18), DS2.ORANGE)
add_text(s, DS2.SLIDE_W - Inches(2.2), Inches(0.85), Inches(1.6), Inches(0.35),
         "TWF2026 PORTAL", size=DS2.SIZE_CAPTION, color=DS2.WHITE, bold=True,
         font=DS2.FONT_HEAD, align=PP_ALIGN.LEFT)
# Kicker
add_text(s, Inches(1.2), Inches(2.1), DS2.SLIDE_W - Inches(2.4), Inches(0.4),
         "EXECUTIVE BRIEFING  —  2026 / 5 / 21",
         size=DS2.SIZE_CAPTION, color=DS2.ORANGE, bold=True, font=DS2.FONT_HEAD)
# Title (大、navy 背景に白)
add_text(s, Inches(1.2), Inches(2.65), DS2.SLIDE_W - Inches(2.4), Inches(2.3),
         "TWF2026\n生産性向上ソリューションコーナー",
         size=DS2.SIZE_DISPLAY - 8, color=DS2.WHITE, bold=True, font=DS2.FONT_HEAD,
         line_spacing=1.1)
add_text(s, Inches(1.2), Inches(5.05), DS2.SLIDE_W - Inches(2.4), Inches(0.6),
         "人手不足対策の本命企画 — 11 社一括提案",
         size=DS2.SIZE_SUBHEAD, color=DS2.ORANGE, bold=True, font=DS2.FONT_HEAD)
# 区切り横線
add_rect(s, Inches(1.2), Inches(5.85), Inches(0.6), Pt(2), DS2.ORANGE)
# 宛先と提案者
add_text(s, Inches(1.2), Inches(6.05), DS2.SLIDE_W - Inches(2.4), Inches(0.35),
         "PREPARED FOR",
         size=DS2.SIZE_TINY, color=DS2.ORANGE, bold=True, font=DS2.FONT_HEAD)
add_text(s, Inches(1.2), Inches(6.35), Inches(7), Inches(0.4),
         "コイケ酸商株式会社 御中",
         size=DS2.SIZE_SUBHEAD - 2, color=DS2.WHITE, bold=True, font=DS2.FONT_HEAD)
add_text(s, Inches(1.2), Inches(6.85), DS2.SLIDE_W - Inches(2.4), Inches(0.4),
         "マツモト産業株式会社 京葉営業所  柏原",
         size=DS2.SIZE_CAPTION, color=DS2.WHITE, font=DS2.FONT_BODY)
set_notes(s, """本日はお時間をいただきありがとうございます。マツモト産業京葉営業所の柏原です。
6 月 12-13 日に幕張メッセで開催する TWF2026「生産性向上ソリューションコーナー」のご案内です。
本日のテーマは、コイケ酸商様の経営層・営業マンが、お客様の人手不足対策をご提案できる「商談装置」として
TWF2026 をご活用いただく、というご提案です。""")


# ============ Slide 02: アジェンダ (Light) ============
s = new_slide(prs)
title_header(s, 2, TOTAL, "本日のアジェンダ", kicker="AGENDA")
# 4 つのセクションを大きめのカードで
sections = [
    ("01", "経営層向け結論 + ROI",   "搬送・溶接・塗装・安全・教示・可視化の 6 工程を 11 社で一括提案。即数字でツカむ。"),
    ("02", "11 社の個別訴求",         "ファナック、ダイヘン、フロニウス、ロボットバンク、メサック等を 1 社 1 枚で深掘り。"),
    ("03", "コイケ視点 + portal 活用","コイケ取扱商品との重なり、営業マンの活用場面、portal の普段使いノウハウ。"),
    ("04", "アクションプラン",         "5/21 (本日) から 6/13 (TWF2026 当日) までの動き方を 5 段階で整理。"),
]
y0 = Inches(2.0); row_h = Inches(1.05); gap = Inches(0.15)
for i, (num, head, body) in enumerate(sections):
    top = y0 + i * (row_h + gap)
    # 番号は大きく orange
    add_text(s, Inches(0.6), top + Inches(0.15), Inches(1.5), Inches(0.85),
             num, size=DS2.SIZE_STAT_M, color=DS2.ORANGE, bold=True,
             font=DS2.FONT_HEAD, align=PP_ALIGN.LEFT, line_spacing=1.0)
    # タイトル
    add_text(s, Inches(2.2), top + Inches(0.15), Inches(4.8), Inches(0.4),
             head, size=DS2.SIZE_SECTION, color=DS2.NAVY_DARK, bold=True, font=DS2.FONT_HEAD)
    # 説明
    add_text(s, Inches(2.2), top + Inches(0.6), DS2.SLIDE_W - Inches(2.8), Inches(0.5),
             body, size=DS2.SIZE_BODY, color=DS2.GRAY_TEXT)
    # 区切り (薄)
    if i < len(sections) - 1:
        add_rect(s, Inches(0.6), top + row_h, DS2.SLIDE_W - Inches(1.2), Pt(0.5), DS2.GRAY_LINE)
footer(s)
set_notes(s, """本日のアジェンダは 4 部構成です。最初に経営層向けの結論と ROI 数字でツカみ、
そこから 11 社の個別訴求を 1 社 1 枚で見ていきます。後半はコイケ酸商様の視点に
寄せたお話 (営業マン活用場面、portal の使い方) と、最後に本日からの動き方を
5 段階で整理します。""")


# ============ Slide 03: ROI 一覧 (★最重要、大数字 72pt) ============
s = new_slide(prs)
title_header(s, 3, TOTAL, "経営層が一番好きな数字", kicker="KEY ROI METRICS")

card_w = Inches(2.9); card_h = Inches(4.5); gap = Inches(0.15)
rois = [
    ("200%", "搬送量向上",       "食品工場 AMR 導入事例",            "ロボットバンク StarLift"),
    ("90%",  "教示時間削減",     "VCOLP 採用、22 メーカー対応",       "ゼネテック OLP"),
    ("47-64%", "塗料使用量削減",  "自動車部品モリブデン / 光輝塗装", "メサック G05 / G07"),
    ("1/4",  "作業者負担",       "元古鉄工事例 4 人/週 → 1 人/日",   "オプティレーザー"),
]
y0 = Inches(2.0)
for i, (num, lab, ev, src) in enumerate(rois):
    left = Inches(0.6 + i * (card_w.inches + gap.inches))
    stat_callout(s, left, y0, card_w, card_h, num, lab, ev, src)
# 下部 callout
add_text(s, Inches(0.6), Inches(6.7), DS2.SLIDE_W - Inches(1.2), Inches(0.4),
         "いずれもメーカー公表値・採用事例ベース — 客先で即提案可能",
         size=DS2.SIZE_BODY, color=DS2.NAVY_DARK, bold=True, font=DS2.FONT_HEAD)
footer(s)
set_notes(s, """このスライドが本日一番大事です。搬送 200%、教示 90%、塗料 47-64%、作業者負担 1/4。
4 つとも公表値ベースの数字なので、客先で「メーカー公表値です」と即ご提案いただけます。
特に塗料削減 (メサック) と元古鉄工事例 4 倍効率 (オプティ) は経営層に強く刺さります。
コイケ営業マンが「導入したらいくら下がりますか」と聞かれた時、この数字でツカんで
個別商談に持ち込む流れが理想です。""")


# ============ Slide 04: 6 工程 × 11 社 マトリクス ============
s = new_slide(prs)
title_header(s, 4, TOTAL, "6 工程 × 11 社の一括提案", kicker="PROCESS COVERAGE MATRIX")
# 6 つの工程を 2x3 で配置
processes = [
    ("01", "🤖", "搬送",       "ロボットバンク / シンテック"),
    ("02", "🔥", "溶接",       "ファナック / ダイヘン / フロニウス\nノビテック / オートスイング"),
    ("03", "🎨", "塗装",       "メサック"),
    ("04", "🛡", "安全",       "小森安全機研究所"),
    ("05", "📐", "教示",       "ゼネテック"),
    ("06", "✨", "クリーニング", "オプティレーザーソリューションズ"),
]
y0 = Inches(1.9); card_w = Inches(4.0); card_h = Inches(2.45); gap_x = Inches(0.1); gap_y = Inches(0.15)
for i, (num, icon, head, makers) in enumerate(processes):
    col = i % 3; row = i // 3
    left = Inches(0.6 + col * (card_w.inches + gap_x.inches))
    top = y0 + row * (card_h + gap_y)
    # 白カード + 上オレンジ細帯
    add_rect(s, left, top, card_w, card_h, DS2.WHITE, line_color=DS2.GRAY_LINE)
    add_rect(s, left, top, card_w, Inches(0.12), DS2.ORANGE)
    # 番号 (orange、small)
    add_text(s, left + Inches(0.25), top + Inches(0.25), Inches(0.5), Inches(0.3),
             num, size=DS2.SIZE_CAPTION, color=DS2.ORANGE_DARK, bold=True, font=DS2.FONT_HEAD)
    # アイコン (大)
    add_text(s, left + card_w - Inches(0.8), top + Inches(0.25), Inches(0.6), Inches(0.6),
             icon, size=24, color=DS2.NAVY_DARK, align=PP_ALIGN.RIGHT)
    # 見出し
    add_text(s, left + Inches(0.25), top + Inches(0.7), card_w - Inches(0.5), Inches(0.55),
             head, size=DS2.SIZE_SECTION + 4, color=DS2.NAVY_DARK, bold=True, font=DS2.FONT_HEAD)
    # メーカー
    add_text(s, left + Inches(0.25), top + Inches(1.35), card_w - Inches(0.5), Inches(1.0),
             makers, size=DS2.SIZE_BODY - 1, color=DS2.GRAY_TEXT, line_spacing=1.35)
footer(s)
set_notes(s, """6 つの工程を 1 つのスライドで俯瞰します。「人手不足対策」は単一の工程の問題ではなく、
搬送・溶接・塗装・安全・教示・クリーニングの 6 工程それぞれにあります。
コイケ酸商様のお客様が「うちは溶接の課題が…」と相談された時、
TWF の溶接 5 社の中から最適な 1-2 社を即提案できる構造になっています。""")


# ============ Slide 05: コイケ取扱 × TWF メーカーの重なり ============
s = new_slide(prs)
title_header(s, 5, TOTAL, "コイケ商材と TWF メーカーの重なり", kicker="OVERLAP MAP")
overlap = [
    ("溶接材料 / ワイヤ", "フロニウス Fortis / ダイヘン VC8 / ファナック CRX / ノビテック Cavitar"),
    ("ガス・ガス機器",    "メサック G05 / G07 / G08 自動ガン、塗装ブース (空気消費 35-400NL/min)"),
    ("保護具・安全用品",  "オートスイング WG3+ ヘルメット、小森 SRD 3D レーダー"),
    ("研磨・砥石",        "ファナック協働ロボパッケージ (ハンドチェンジャーでグラインダー対応)"),
    ("搬送機器",          "ロボットバンク StarLift / Star-7、シンテック T-Arm / Rail Station"),
    ("自動化・ティーチング", "ゼネテック VCOLP (22 メーカー対応)、オプティレーザー (下地処理)"),
]
y0 = Inches(2.05); row_h = Inches(0.7); gap = Inches(0.05)
for i, (left_col, right_col) in enumerate(overlap):
    top = y0 + i * (row_h + gap)
    # 左カラム
    add_rect(s, Inches(0.6), top, Inches(3.4), row_h, DS2.GRAY_BG, line_color=DS2.GRAY_LINE)
    add_text(s, Inches(0.75), top + Inches(0.18), Inches(3.1), Inches(0.4),
             left_col, size=DS2.SIZE_BODY, color=DS2.NAVY_DARK, bold=True, font=DS2.FONT_HEAD)
    # 矢印 (orange small rectangle)
    add_text(s, Inches(4.05), top + Inches(0.18), Inches(0.4), Inches(0.4),
             "→", size=DS2.SIZE_SECTION, color=DS2.ORANGE_DARK, bold=True, align=PP_ALIGN.CENTER)
    # 右カラム
    add_rect(s, Inches(4.5), top, DS2.SLIDE_W - Inches(5.1), row_h, DS2.WHITE, line_color=DS2.GRAY_LINE)
    add_text(s, Inches(4.65), top + Inches(0.18), DS2.SLIDE_W - Inches(5.4), Inches(0.4),
             right_col, size=DS2.SIZE_BODY, color=DS2.GRAY_TEXT)
# 締めのメッセージ
add_text(s, Inches(0.6), Inches(6.65), DS2.SLIDE_W - Inches(1.2), Inches(0.45),
         "コイケのお客様の現場課題は、TWF メーカー製品で生産性向上に直結する。",
         size=DS2.SIZE_BODY, color=DS2.ORANGE_DARK, bold=True, font=DS2.FONT_HEAD)
footer(s)
set_notes(s, """コイケ酸商様の主力 6 商材と、TWF 出展 11 社の関係を整理しました。
従来のコイケ商材 (溶接材料・ガス・保護具) と最近の課題対応商材 (自動化・搬送・教示) の
橋渡しが TWF メーカーで一気に揃います。「コイケのお客様で人手不足の現場があれば、
TWF のこの 11 社で答えが出る」という売り方ができます。""")


# ============ Slide 06-17: メーカー 11 社個別 (12 枚、内 4 枚動画) ============

def maker_slide(prs, page_num, maker, tagline, message, points, materials, portal_url, notes, video_label=None):
    s = new_slide(prs)
    title_header(s, page_num, TOTAL, maker, kicker=tagline)
    # 左: 経営層メッセージカード (light gray bg + orange top stripe)
    msg_top = Inches(2.0); msg_h = Inches(2.8); msg_w = Inches(5.0)
    add_rect(s, Inches(0.6), msg_top, msg_w, Inches(0.12), DS2.ORANGE)
    add_rect(s, Inches(0.6), msg_top + Inches(0.12), msg_w, msg_h - Inches(0.12), DS2.GRAY_BG)
    add_text(s, Inches(0.85), msg_top + Inches(0.3), msg_w - Inches(0.5), Inches(0.35),
             "EXECUTIVE MESSAGE", size=DS2.SIZE_TINY, color=DS2.ORANGE_DARK, bold=True, font=DS2.FONT_HEAD)
    add_text(s, Inches(0.85), msg_top + Inches(0.75), msg_w - Inches(0.5), msg_h - Inches(1.0),
             message, size=DS2.SIZE_BODY, color=DS2.NAVY_DARK, bold=True, line_spacing=1.5)
    # 右: TWF みどころ bullets
    pts_left = Inches(5.9); pts_w = DS2.SLIDE_W - pts_left - Inches(0.6)
    add_text(s, pts_left, msg_top + Inches(0.18), pts_w, Inches(0.35),
             "TWF2026 MIDOKORO", size=DS2.SIZE_TINY, color=DS2.ORANGE_DARK, bold=True, font=DS2.FONT_HEAD)
    add_bullets(s, pts_left, msg_top + Inches(0.65), pts_w, Inches(3.0), points, size=DS2.SIZE_BODY - 1)
    # 動画 (該当時)
    if video_label:
        video_frame(s, Inches(0.6), Inches(5.0), Inches(7.5), Inches(1.95), video_label)
    # 下部 thin info bar (材料 + portal URL)
    info_top = DS2.SLIDE_H - Inches(1.1)
    add_rect(s, Inches(0.6), info_top, DS2.SLIDE_W - Inches(1.2), Pt(0.5), DS2.GRAY_LINE)
    add_text(s, Inches(0.6), info_top + Inches(0.15), DS2.SLIDE_W - Inches(1.2), Inches(0.3),
             f"📄  配布資料: {materials}", size=DS2.SIZE_CAPTION, color=DS2.GRAY_TEXT)
    add_text(s, Inches(0.6), info_top + Inches(0.5), DS2.SLIDE_W - Inches(1.2), Inches(0.3),
             f"🌐  詳細ページ: {portal_url}",
             size=DS2.SIZE_CAPTION, color=DS2.ORANGE_DARK, bold=True)
    footer(s)
    set_notes(s, notes)


maker_slide(prs, 6, "ファナック㈱",
    "FANUC | 協働ロボのマルチタスク化 — 3kg 可搬 × 協働ロボパッケージ × TIG フィラー",
    "省スペース工場で 1 台を溶接・研磨・TIG に使い回すマルチタスク化。安全柵不要前提で固定設備化せず、小型協働ロボで小さく始められる。",
    [
        "🆕 3kg 可搬 マグネット式 高電圧タッチセンサー、軽量 11kg",
        "🤖 協働ロボパッケージ 安全柵不要前提、狭い工場対応",
        "🔁 ワンタッチハンドチェンジャー CO2 トーチ↔グラインダー / TIG-研磨 / ハンドレーザ",
        "💪 力覚研磨 マツモト機械フローティング + 内蔵力覚センサ",
        "🔥 TIG フィラー仕様 簡単教示で高品質・高能率パルス TIG",
        "🎯 連携実績 大阪・神奈川 WF で ATC / HW1000 等を実演、TWF2026 でも継続",
    ],
    "TWF2026 公式パネル (2 ページ統合) + 3kg 可搬仕様 PDF",
    "twf2026-portal.pages.dev/m/fanuc/",
    """ファナックは協働ロボのマルチタスク化を 4 つの切り口で提案します。
3kg 可搬は工場の好きな場所にマグネットでくっつけて使える「持ち運べる協働ロボ」、
協働ロボパッケージは安全柵なしでハンドチェンジャー付き、TIG フィラー仕様は熟練 TIG 作業を自動化、
そして過去 WF で見せてきた周辺メーカー連携も TWF2026 で継続展示します。""")

maker_slide(prs, 7, "㈱ダイヘン",
    "DAIHEN | TIG/MAG 兼用仕様 + AiTran 連携自動化 — 溶接品質の安定化",
    "1 台架台に CO2・TIG 溶接機を搭載、段取り替え簡単。ダイヘン独自制御で高軌跡精度、難易度の高い TIG フィラーでも安定高品質。",
    [
        "🔁 TIG/MAG 兼用仕様 1 台で 2 方式、段取り替えで使い分け",
        "🎯 高軌跡精度 TIG フィラー溶接も安定高品質",
        "🤖 AiTran 連携 搬送→位置補正→溶接の一気通貫自動化",
        "🚀 マツモト産業ブース連動 自動化推進コーナーで実機デモ",
        "💡 Slide 18 で「実演セミナーで同じデモが見られる」と連動訴求",
    ],
    "TWF2026 公式パネル + 実演セミナー (Slide 18 で詳細)",
    "twf2026-portal.pages.dev/m/daihen/",
    """ダイヘンの目玉は VC8 + AiTran 連携デモです。動画は神奈川 WF で撮影した実機デモを
そのまま埋め込みます。AiTran が部材を運んで、VC8 が位置補正して溶接する、という
一連の自動化が 2 分半でわかります。Slide 18 (実演セミナー) で同じ内容の生実演があるため、
「動画 → 生で見たい」の動線で来場動機を作ります。""",
    video_label="ダイヘン VC8 × AiTran500 連携デモ (2:29) — YouTube -ydKdIio5es")

maker_slide(prs, 8, "フロニウスジャパン㈱",
    "FRONIUS | Fortis シリーズ 270〜500A — Wizard 機能で若手でも条件設定",
    "MIG/MAG・TIG・手棒・ガウジングまで 1 シリーズで対応、Wizard 機能で経験浅でも溶接条件設定が可能、人材育成にも有効。",
    [
        "🆕 Fortis シリーズ 270〜500A、空冷/水冷、送給装置一体型/別置き型",
        "🧙 Wizard 機能 経験浅でも溶接条件設定が可能",
        "🔥 幅広い工法 MIG/MAG 直流・パルス、TIG、手棒、ガウジング",
        "🤖 自動化連携 ファナック CRX + TPS500i + CMT の協働ロボ連携",
        "🎁 TWF2026 ご注文特典 自動遮光面 Vizor 4000 Plus プレゼント",
    ],
    "TWF2026 公式パネル + Manual 溶接機チラシ",
    "twf2026-portal.pages.dev/m/furoniusujapan/",
    """フロニウスの強みは Wizard 機能です。若手作業者でも条件設定が間違いなくできるので、
人材育成と品質安定の両立に効きます。経験豊富なベテランが減る中で「機械側で知見を担保する」
という発想は経営層に強く刺さります。
注文特典は現場担当者向け、経営層プレゼンでは Wizard が主訴求です。""")

maker_slide(prs, 9, "ロボットバンク㈱",
    "ROBOTBANK | StarLift + Star-7 — 搬送と清掃の無人化、5 業界事例",
    "AMR の 5 系統ラインナップで「人が歩く仕事」を置き換え。食品工場で搬送量 200% 向上、修理工場で操作教育 30 分で即戦力化。",
    [
        "🤖 StarLift 150 / 300 / 600 積載 150〜600kg をカバー",
        "🌐 全シリーズ 5 系統 StarShip / StarMax / StarLight (低床 31cm 棚下くぐり) / RisuBot",
        "🧹 Star-7 業務用清掃ロボット 拭き取り・掃除・吸塵・磨き (公式パネル新情報)",
        "🎯 動作仕様 段差 20mm、登坂 8°、最小通過幅 60cm、稼働 10h、±3cm 精度",
        "👀 採用事例 食品 200% / 修理 30 分教育 / 部品製造 導入 2 日 / 自動車部品 24 時間体制",
    ],
    "TWF2026 公式パネル + 導入事例集 + 製品ハイライト + 搬送ロボットカタログ",
    "twf2026-portal.pages.dev/m/robottobanku/",
    """ロボットバンクは AMR 搬送ロボの 5 系統と、新発表の Star-7 清掃ロボットの 2 本柱です。
「人が歩く仕事を置き換える」という単一コンセプトで覚えていただけるのがポイントです。
食品工場での 200% 搬送量向上の事例は、コイケ酸商様のお客様の中にも食品・部品製造の
現場があれば、直接適用しやすい数字です。""")

maker_slide(prs, 10, "㈱メサック",
    "MESAC | ロボットつかみ方式塗装ブース — 1㎡省スペース × 塗料 47〜64% 削減",
    "ガンを固定しロボットがワークを持つ逆転発想で、塗装ブース自体が約 1㎡。G05 自動ガンで塗料使用量 47〜64% 削減事例 (自動車部品)。",
    [
        "📐 設置面積 約 1 ㎡で省スペース設置",
        "💨 排気風量 30 ㎥/min、ガン下向きで塗料飛散範囲限定",
        "🎯 塗料供給経路 ポンプ〜ガン間ホース 約 1m で短縮化",
        "🆕 G05/G07/G08 自動ガン 塗料使用量 47-64% 節約事例、G08 ダイヤフラム構造 2 液塗料対応",
        "🧹 清掃性 ガン下向き固定、ロールフィルター切断式メンテで容易化",
        "🔧 工具レス分解 G08 で塗料ブロックのみ取り外し可",
    ],
    "TWF2026 公式パネル + G05/G07/G08 PDF + コンパクト塗装ブースちらし",
    "twf2026-portal.pages.dev/m/mesakku/",
    """メサックは本日の ROI スライド (Slide 3) で出した「塗料 47-64% 削減」の出元です。
塗料は塗装現場のコストの大半を占めるので、そのまま経費削減として試算しやすい数字です。
さらに塗装ブース自体が 1㎡で済むという省スペース化と、排気風量 30㎥/min で空調コスト
も下がるので、トリプルでコスト効果があります。""")

maker_slide(prs, 11, "㈱ゼネテック",
    "GENETECH | Visual Components Robotics OLP — ティーチング 90% 削減 × 22 メーカー対応",
    "ワーク 3D-CAD クリックで 6 軸ロボット動作軌跡を作成。22 メーカーのロボットプログラム出力対応、複数メーカー混在ラインでも標準化可能。",
    [
        "🎯 教示時間 90% 削減 VCOLP 採用で 1/10",
        "🤖 22 メーカー対応 ロボットプログラム出力対象",
        "📐 CAD クリック教示 ワーク 3D-CAD クリックで 6 軸ロボット動作軌跡",
        "🆕 VCOLP 5.0 2026 年 3 月 18 日提供開始",
        "👀 5 用途別画面 アーク/スポット/切断/研磨/塗装",
    ],
    "TWF2026 公式パネル + VCOLP パンフレット",
    "twf2026-portal.pages.dev/m/zenetekku/",
    """ゼネテックの VCOLP は教示時間 90% 削減という強い数字です。ロボット教示は熟練者が
時間をかけて行う作業で、人件費と機会損失の両方の負担があります。それが 1/10 になる
というのは、ロボット導入の意思決定を後押しする数字です。
22 メーカー対応なので、コイケ酸商様のお客様がどのメーカーのロボットを持っていても
提案できる汎用性があります。""")

maker_slide(prs, 12, "㈱小森安全機研究所",
    "KOMORI | SRD 3D レーダー安全システム — 世界初 SIL2/PLd 規格準拠",
    "光・粉塵・煙・水・雨に強い 3D 安全レーダー。安全投資 = 経営価値 (労災ゼロ、設備停止リスク低減、ライン稼働率向上)。",
    [
        "🛡 世界初 SIL2/PLd 規格準拠 3D 安全レーダー SRD シリーズ",
        "📡 SRD 仕様 60GHz FMCW、応答 100ms 以下、最大 6 センサ、検知距離 5m/9m/4m",
        "🌧 耐環境性 光・粉塵・煙・水・雨に強く、降雨量 45mm/h 対応",
        "🎯 動的検知ゾーン 検知ゾーン/警告ゾーンを動的設定、最大 32 種類",
        "👀 光学式が苦手な外乱環境 (溶接・塗装・搬送・屋外設備) で相談可能",
    ],
    "TWF2026 公式パネル + SRD + AI カメラ KAG 製品案内",
    "twf2026-portal.pages.dev/m/komori-anzen-ki-kenkyuusho/",
    """小森安全機の SRD は世界初の SIL2/PLd 規格準拠 3D レーダーです。
光学センサーが苦手な現場 (粉塵・煙・水・雨) でも安全検知できるので、
労災ゼロ + 設備停止リスク低減 = ライン稼働率向上、というロジックでご提案ください。
経営層は安全を「コスト」と捉えがちですが、「稼働率」と言い換えると意思決定が早くなります。""")

maker_slide(prs, 13, "シンテック㈱",
    "SHINTECH | 3arm / T-Arm / Rail Station — 作業負荷 × 労災リスク低減",
    "重量物・工具保持・搬送補助の作業負荷を下げ、腰痛・労災リスクも低減。トヨタ 6000+/日野 3000+/ダイハツ 500+ セットの自動車メーカー採用実績。",
    [
        "💪 3arm 締付・組立・バリ取り・持上操作、最大荷重 35kg",
        "🏗 T-Arm 耐荷重 40〜650kg、オートバランス標準装備",
        "📈 強度向上 引張・圧縮強度 1.8 倍に再設計",
        "🏭 採用実績 トヨタ 6000+ / 日野 3000+ / ダイハツ 500+ セット",
        "🛤 Rail Station 落下防止に優れた運搬搭載補助レール",
    ],
    "TWF2026 公式パネル + 3arm カタログ + 製品プレゼン",
    "twf2026-portal.pages.dev/m/shintech/",
    """シンテックは作業負荷低減の老舗です。トヨタ 6000、日野 3000、ダイハツ 500 セット
の採用実績は信頼性の証です。「ベテラン作業者の離職リスク」「腰痛による労災コスト」
「採用しても定着しない」を 3 製品の組合せで解決できます。
省人化ではなく「働きやすさで定着率向上」というニュアンスでご提案ください。""")

maker_slide(prs, 14, "㈱ノビテック",
    "NOVITEC | Cavitar Welding Camera + Weld-Eye — 溶接不良リアルタイム可視化",
    "アーク光・ヒュームをカットし溶融池を可視化、Weld-Eye で溶融池・面積・溶融速度を AI リアルタイム判定。",
    [
        "👀 溶融池可視化 アーク光・ヒュームをカットし詳細観察",
        "🤖 AI 解析 Weld-Eye で溶融池・面積・溶融速度を AI 解析、不良可否リアルタイム判定",
        "🎯 2 機種比較 C350 / C400 のサイズ・fps・設置性",
        "📐 高解像度 最大 1,440×1,080 px、C350 は最大 500 fps",
        "👀 高品位ワーク要求現場 (品質保証・原因究明・教育用途) 向け",
    ],
    "TWF2026 公式パネル + Cavitar カタログ",
    "twf2026-portal.pages.dev/m/nobitekku/",
    """ノビテックの Cavitar Welding Camera は、アーク光に邪魔されずに溶接プールを
直接観察できるカメラです。Weld-Eye と組み合わせると AI が溶接の良否をリアルタイム判定
できるので、品質保証の証跡として動画記録できます。
高品位ワーク (航空・自動車・医療機器) のお客様で「品質エビデンスを残したい」
というニーズがあれば即提案できます。""")

maker_slide(prs, 15, "㈱オートスイング (OTOS) — Ray-X",
    "OTOS #1 | WGC200 / WGC400 溶接カメラ — アーク溶接そのものが鮮明に見える",
    "補助照明不要・超小型で、ロボット溶接やキャリッジ溶接にも設置しやすい。品質確認・技能教育・遠隔監視用途を具体化できる。",
    [
        "👀 Ray-X アーク溶接そのものが鮮明、補助照明不要",
        "📷 2 機種 WGC200 (有線、152g) / WGC400 (Wi-Fi6 無線、274g)",
        "📐 計測オプション アーク長・ビード幅・シーム位置",
        "🎬 サンプル動画 GMAW / FCAW / Orbital GTAW 撮影",
    ],
    "TWF2026 公式パネル + OTOSWING リーフレット",
    "twf2026-portal.pages.dev/m/ootosuingu-otos/",
    """OTOS の Ray-X は超小型の溶接カメラで、ロボット溶接にもキャリッジ溶接にも
設置しやすいのが特徴です。動画はこの後 Slide 16 と合わせて 2 種類お見せします。
本スライドは Ray-X の代表的な撮影サンプルです。""",
    video_label="OTOS Ray-X 撮影サンプル動画 #1 (柏原所有 MP4)")

maker_slide(prs, 16, "㈱オートスイング (OTOS) — WG3+",
    "OTOS #2 | カメラ搭載溶接ヘルメット — 品質確認 × 技能伝承",
    "作業者目線の可視化で、熟練者のノウハウを動画として記録・共有可能。後継者教育、品質トレーニング、遠隔指導に活用。",
    [
        "🪖 WG3+ ヘルメット カメラ搭載溶接ヘルメットで作業者目線の可視化",
        "🎯 技能伝承 熟練者のノウハウを動画で記録・共有",
        "👀 品質確認 溶接プール・シーム位置を作業者視点で確認",
        "🎬 教育用途 後継者教育、品質トレーニングに活用",
    ],
    "TWF2026 公式パネル + 関連製品リーフレット",
    "twf2026-portal.pages.dev/m/ootosuingu-otos/",
    """WG3+ はカメラがついた溶接ヘルメットです。「熟練者がどう見ているか」を
そのまま動画で残せるので、後継者育成や品質基準の標準化に使えます。
コイケ酸商様のお客様の中で「熟練者が引退する前にノウハウを残したい」というニーズが
あれば、保護具予算と教育予算の両方から正当化できます。""",
    video_label="OTOS WG3+ ヘルメット撮影サンプル動画 #2 (柏原所有 MP4)")

maker_slide(prs, 17, "オプティレーザーソリューションズ㈱",
    "OPTILASER | ULT LASER シリーズ — 元古鉄工事例で作業者負担 1/4",
    "錆・塗膜・酸化皮膜を非接触で除去する国内生産レーザークリーナー。元古鉄工事例で 4 人/1 週間 → 1 人/1 日に作業時間短縮。",
    [
        "🏭 元古鉄工事例 4 人/1 週間 → 1 人/1 日、作業者負担 1/4",
        "🇯🇵 国内生産 大阪本社のレーザークリーナー専門メーカー",
        "🔦 9 機種ラインナップ CW / Pulse の用途別",
        "⚡ 即起動 最短 5 秒、空冷 6h 以上・水冷 10h 以上の連続照射",
        "✨ 環境負荷低 薬品・研磨材を使わず二次廃棄物を抑える",
    ],
    "TWF2026 公式パネル + 公式カタログ + 元古鉄工事例記事",
    "twf2026-portal.pages.dev/m/oputeireezaasoryuushonzu/",
    """オプティレーザーは初 TWF 出展の注目株です。元古鉄工事例は ROI スライドでも
出した「4 人/1 週間 → 1 人/1 日」という強い数字の出元です。
作業者負担 4 分の 1 + 工期 1/5 + 二次廃棄物ほぼゼロという、3 重の効果があります。
鉄骨ファブ、橋梁、船舶など重厚長大系のお客様に即訴求できます。""",
    video_label="オプティレーザー 製品紹介 (YouTube ypxAtVayQxQ)")


# ============ Slide 18: 実演セミナー紹介 ============
s = new_slide(prs)
title_header(s, 18, TOTAL, "実演セミナー 4 社 — 目の前で動く実機を体感", kicker="LIVE DEMOS")
seminars = [
    ("3M",         "🪜", "フルハーネスデモンストレーション (墜落防止実験 + 吊り下げ体験)"),
    ("神戸製鋼所", "🔥", "AXELARC™ 新ワイヤ送給制御プロセス、建機向け安定アーク × 低スパッタ"),
    ("ダイヘン",   "🤖", "VC8 + AiTran 連携デモ — Slide 07 で見せた動画と同内容の生実演!"),
    ("三菱電機",   "✨", "二次元ファイバレーザ加工機 ML3015GX-F60、高速・高精度切断"),
]
y0 = Inches(2.0); row_h = Inches(1.05); gap = Inches(0.15)
for i, (name, icon, desc) in enumerate(seminars):
    top = y0 + i * (row_h + gap)
    # アイコンサークル (orange tint) + name
    add_rect(s, Inches(0.6), top, DS2.SLIDE_W - Inches(1.2), row_h, DS2.WHITE, line_color=DS2.GRAY_LINE)
    # オレンジ左端ストリップ
    add_rect(s, Inches(0.6), top, Inches(0.15), row_h, DS2.ORANGE)
    add_text(s, Inches(0.9), top + Inches(0.15), Inches(0.8), Inches(0.75),
             icon, size=32, color=DS2.NAVY_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.85), top + Inches(0.15), Inches(3.2), Inches(0.4),
             name, size=DS2.SIZE_SECTION - 2, color=DS2.NAVY_DARK, bold=True, font=DS2.FONT_HEAD)
    add_text(s, Inches(1.85), top + Inches(0.55), DS2.SLIDE_W - Inches(2.6), Inches(0.5),
             desc, size=DS2.SIZE_BODY - 1, color=DS2.GRAY_TEXT)
# Callout
add_rect(s, Inches(0.6), Inches(6.9), DS2.SLIDE_W - Inches(1.2), Inches(0.45), DS2.NAVY_DARK)
add_text(s, Inches(0.85), Inches(6.95), DS2.SLIDE_W - Inches(1.7), Inches(0.35),
         "Slide 07 のダイヘン動画 → 実演セミナーで生実演、という動線で来場動機を作る",
         size=DS2.SIZE_BODY, color=DS2.ORANGE, bold=True, font=DS2.FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)
footer(s)
set_notes(s, """実演セミナーは目の前で動く実機を体感できる無料セミナーです。4 社開催で、
特にダイヘンの VC8 + AiTran 連携デモは、本日 Slide 07 でお見せした動画と
同じものが生実演されます。「動画で見た → 実機で確かめたい」という動機付けで、
コイケ酸商様のお客様の来場誘導に効きます。""")


# ============ Slide 19: 作業環境向上 + 初TWF出展 (2 column) ============
s = new_slide(prs)
title_header(s, 19, TOTAL, "周辺企画 — 作業環境向上 + 初TWF出展", kicker="ADJACENT VALUE")
# Left card
left_w = Inches(6.0); card_h = Inches(4.7)
add_rect(s, Inches(0.6), Inches(2.0), left_w, Inches(0.12), DS2.ORANGE)
add_rect(s, Inches(0.6), Inches(2.12), left_w, card_h - Inches(0.12), DS2.WHITE, line_color=DS2.GRAY_LINE)
add_text(s, Inches(0.85), Inches(2.3), left_w - Inches(0.5), Inches(0.45),
         "作業環境向上ブース", size=DS2.SIZE_SECTION, color=DS2.NAVY_DARK, bold=True, font=DS2.FONT_HEAD)
add_bullets(s, Inches(0.85), Inches(2.85), left_w - Inches(0.5), Inches(3.5),
            [
                "🌬 集塵・脱臭設備の最新提案",
                "🦺 防塵マスク・遮光面・フルハーネス等の安全衛生製品",
                "🔇 騒音対策、空調機器、健康診断ブース",
                "💡 「働きやすい現場」= 「採用しやすい現場」へつなぐ訴求",
                "🏭 コイケのお客様の現場改善提案にも直結",
            ], size=DS2.SIZE_BODY)
# Right card
right_left = Inches(0.6) + left_w + Inches(0.15); right_w = DS2.SLIDE_W - right_left - Inches(0.6)
add_rect(s, right_left, Inches(2.0), right_w, Inches(0.12), DS2.NAVY_DARK)
add_rect(s, right_left, Inches(2.12), right_w, card_h - Inches(0.12), DS2.WHITE, line_color=DS2.GRAY_LINE)
add_text(s, right_left + Inches(0.25), Inches(2.3), right_w - Inches(0.5), Inches(0.45),
         "初TWF出展いちおしメーカー", size=DS2.SIZE_SECTION, color=DS2.NAVY_DARK, bold=True, font=DS2.FONT_HEAD)
add_bullets(s, right_left + Inches(0.25), Inches(2.85), right_w - Inches(0.5), Inches(3.5),
            [
                "🎯 オプティレーザーソリューションズ (Slide 17)",
                "🎯 ゼネテック VCOLP 5.0 (Slide 11)",
                "🎯 ロボットバンク Star-7 清掃 (Slide 09)",
                "🎯 他、TWF2026 で初お披露目メーカー多数",
                "💡 業界トレンドの先取り提案として活用可能",
            ], size=DS2.SIZE_BODY)
footer(s)
set_notes(s, """作業環境向上ブースは「働きやすい現場 = 採用しやすい現場」という切り口で
ご紹介ください。初 TWF 出展のメーカー (オプティ・ゼネテック VCOLP 5.0・Star-7) は
業界トレンドの先取り情報として、コイケ営業マンが「他社より早い情報」を客先に届けられます。""")


# ============ Slide 20: portal 紹介 + QR ============
s = new_slide(prs)
title_header(s, 20, TOTAL, "TWF2026 みどころポータル — 営業マンの普段使い資料", kicker="THE PORTAL")
# Left: URL + features
left_w = Inches(7.8); top0 = Inches(2.0)
add_rect(s, Inches(0.6), top0, left_w, Inches(0.12), DS2.ORANGE)
add_rect(s, Inches(0.6), top0 + Inches(0.12), left_w, Inches(4.7), DS2.WHITE, line_color=DS2.GRAY_LINE)
add_text(s, Inches(0.85), top0 + Inches(0.3), left_w - Inches(0.5), Inches(0.55),
         "twf2026-portal.pages.dev", size=DS2.SIZE_SECTION + 2, color=DS2.ORANGE_DARK, bold=True, font=DS2.FONT_HEAD)
features = [
    "● 149 出展メーカーの portal、生産性向上 11 社含む全社情報を収録",
    "● 各社の Q1-Q5 (企画概要・新製品・みどころ・特典・配布資料) を整理",
    "● 公式パネル PDF・カタログ・チラシをワンクリックで DL 可能",
    "● スマホ対応、来場前・来場後どちらも快適に閲覧",
    "● 事前来場登録ボタン → 公式特設ページ経由で 1 分登録",
    "● 5/18 時点で生産性向上 11 社の Q1-Q5 が経営層向けに完全リライト済",
]
add_bullets(s, Inches(0.85), top0 + Inches(1.0), left_w - Inches(0.5), Inches(3.5),
            features, size=DS2.SIZE_BODY)
# Right: QR placeholder (navy)
qr_left = Inches(0.6) + left_w + Inches(0.15); qr_w = DS2.SLIDE_W - qr_left - Inches(0.6)
add_rect(s, qr_left, top0, qr_w, Inches(0.12), DS2.NAVY_DARK)
add_rect(s, qr_left, top0 + Inches(0.12), qr_w, Inches(4.7), DS2.NAVY_DARK)
add_text(s, qr_left, top0 + Inches(0.45), qr_w, Inches(0.4),
         "SCAN", size=DS2.SIZE_CAPTION, color=DS2.ORANGE, bold=True, align=PP_ALIGN.CENTER, font=DS2.FONT_HEAD)
# QR code square
qr_box_w = Inches(2.5); qr_box_left = qr_left + (qr_w - qr_box_w) / 2
add_rect(s, qr_box_left, top0 + Inches(1.0), qr_box_w, qr_box_w, DS2.WHITE)
add_text(s, qr_box_left, top0 + Inches(2.05), qr_box_w, Inches(0.5),
         "QR", size=44, color=DS2.NAVY_DARK, bold=True, align=PP_ALIGN.CENTER, font=DS2.FONT_HEAD)
add_text(s, qr_box_left, top0 + Inches(2.7), qr_box_w, Inches(0.4),
         "(配布版でQR挿入)", size=DS2.SIZE_TINY, color=DS2.GRAY_TEXT, align=PP_ALIGN.CENTER)
add_text(s, qr_left, top0 + Inches(4.0), qr_w, Inches(0.4),
         "スマホで即アクセス", size=DS2.SIZE_BODY, color=DS2.WHITE, align=PP_ALIGN.CENTER, font=DS2.FONT_HEAD)
# 締め
add_text(s, Inches(0.6), Inches(6.8), DS2.SLIDE_W - Inches(1.2), Inches(0.45),
         "コイケ営業マンが客先と一緒にスマホで開いて、工程別の見どころをその場で案内できる。",
         size=DS2.SIZE_BODY, color=DS2.NAVY_DARK, bold=True, font=DS2.FONT_HEAD)
footer(s)
set_notes(s, """portal はコイケ営業マンの「普段使い資料」です。客先でスマホを開いて、
お客様と一緒に見ながらご案内する想定で設計してあります。
特に強調したいのは、5/18 時点で生産性向上 11 社の Q1-Q5 が経営層向けに完全リライト済
である点です。営業マンがそのまま読み上げてもプレゼンになるレベルで整理しています。""")


# ============ Slide 21: 資料の使い分け (table-like) ============
s = new_slide(prs)
title_header(s, 21, TOTAL, "資料の使い分け — 場面別の使い方", kicker="MATERIALS PLAYBOOK")
rows = [
    ("公式パネル PDF (10 社分)",   "ブース前説明 / 営業配布用",       "経営層の「掴み」、客先で 30 秒〜1 分で概要伝達。各社 1-2 ページの簡潔な内容、印刷も可。"),
    ("既存パンフレット",           "仕様確認 / 商談深掘り",            "各社 10-30 ページの詳細仕様。技術担当との具体的検討、見積前提条件の整理に使用。"),
    ("3kg 可搬仕様 PDF",            "新製品の単独訴求",                  "ファナック最新協働ロボの 1 ページ簡易資料。デモ依頼受付ありの旨を明示。"),
]
y0 = Inches(2.0); row_h = Inches(1.4); gap = Inches(0.15)
# Header row
add_rect(s, Inches(0.6), y0, Inches(3.5), Inches(0.4), DS2.NAVY_DARK)
add_text(s, Inches(0.75), y0 + Inches(0.06), Inches(3.3), Inches(0.3),
         "資料", size=DS2.SIZE_CAPTION, color=DS2.WHITE, bold=True, font=DS2.FONT_HEAD)
add_rect(s, Inches(4.25), y0, Inches(3.0), Inches(0.4), DS2.NAVY_DARK)
add_text(s, Inches(4.4), y0 + Inches(0.06), Inches(2.8), Inches(0.3),
         "使う場面", size=DS2.SIZE_CAPTION, color=DS2.WHITE, bold=True, font=DS2.FONT_HEAD)
add_rect(s, Inches(7.4), y0, DS2.SLIDE_W - Inches(8.0), Inches(0.4), DS2.NAVY_DARK)
add_text(s, Inches(7.55), y0 + Inches(0.06), DS2.SLIDE_W - Inches(8.2), Inches(0.3),
         "活用ポイント", size=DS2.SIZE_CAPTION, color=DS2.WHITE, bold=True, font=DS2.FONT_HEAD)
for i, (col1, col2, col3) in enumerate(rows):
    top = y0 + Inches(0.5) + i * (row_h + gap)
    add_rect(s, Inches(0.6), top, Inches(3.5), row_h, DS2.GRAY_BG, line_color=DS2.GRAY_LINE)
    add_text(s, Inches(0.75), top + Inches(0.2), Inches(3.3), Inches(0.4),
             col1, size=DS2.SIZE_BODY, color=DS2.NAVY_DARK, bold=True, font=DS2.FONT_HEAD)
    # left orange stripe
    add_rect(s, Inches(0.6), top, Inches(0.1), row_h, DS2.ORANGE)
    add_rect(s, Inches(4.25), top, Inches(3.0), row_h, DS2.WHITE, line_color=DS2.GRAY_LINE)
    add_text(s, Inches(4.4), top + Inches(0.2), Inches(2.8), Inches(1.0),
             col2, size=DS2.SIZE_BODY, color=DS2.ORANGE_DARK, bold=True, font=DS2.FONT_HEAD)
    add_rect(s, Inches(7.4), top, DS2.SLIDE_W - Inches(8.0), row_h, DS2.WHITE, line_color=DS2.GRAY_LINE)
    add_text(s, Inches(7.55), top + Inches(0.2), DS2.SLIDE_W - Inches(8.2), row_h - Inches(0.4),
             col3, size=DS2.SIZE_BODY - 1, color=DS2.GRAY_TEXT, line_spacing=1.35)
footer(s)
set_notes(s, """3 種類の資料の使い分けです。公式パネル PDF は経営層への「掴み」用、
既存パンフレットは技術担当との詳細検討用、3kg 可搬は新製品の単独訴求用。
コイケ営業マンが客先で「まず何を渡すか」迷ったときの判断基準としてご活用ください。""")


# ============ Slide 22: 営業トーク 3 ステップ ============
s = new_slide(prs)
title_header(s, 22, TOTAL, "営業トーク 3 ステップ — 数字 → 体験 → 商談", kicker="SALES SCRIPT")
steps = [
    ("STEP 1", "数字 ROI でツカむ",
     "「塗料 47-64% 削減」「搬送 200% 向上」「教示 90% 削減」を最初に出す。経営層は数字に反応する。"),
    ("STEP 2", "実機体験 / TWF 来場誘導",
     "「実機を目の前で動かしてます、6/12-13 に幕張メッセでぜひご確認を」と来場誘導。公式特設ページから事前登録。"),
    ("STEP 3", "個別相談 / 商談化",
     "「導入規模・タイミング・ご予算感をお聞かせいただければ、メーカーと一緒に最適構成をご提案します」。"),
]
y0 = Inches(2.0); col_w = (DS2.SLIDE_W - Inches(1.5)) / 3; gap = Inches(0.15)
for i, (label, head, body) in enumerate(steps):
    left = Inches(0.6) + i * (col_w + gap)
    # Top dark band
    add_rect(s, left, y0, col_w, Inches(0.7), DS2.NAVY_DARK)
    add_text(s, left, y0 + Inches(0.15), col_w, Inches(0.4),
             label, size=DS2.SIZE_CAPTION + 2, color=DS2.ORANGE, bold=True, align=PP_ALIGN.CENTER, font=DS2.FONT_HEAD)
    # Body
    add_rect(s, left, y0 + Inches(0.7), col_w, Inches(3.5), DS2.WHITE, line_color=DS2.GRAY_LINE)
    add_text(s, left + Inches(0.2), y0 + Inches(0.95), col_w - Inches(0.4), Inches(0.8),
             head, size=DS2.SIZE_SECTION, color=DS2.NAVY_DARK, bold=True, font=DS2.FONT_HEAD, line_spacing=1.3)
    add_text(s, left + Inches(0.2), y0 + Inches(2.0), col_w - Inches(0.4), Inches(2.0),
             body, size=DS2.SIZE_BODY, color=DS2.GRAY_TEXT, line_spacing=1.5)
# Bottom callout
add_rect(s, Inches(0.6), Inches(6.4), DS2.SLIDE_W - Inches(1.2), Inches(0.5), DS2.ORANGE)
add_text(s, Inches(0.85), Inches(6.5), DS2.SLIDE_W - Inches(1.7), Inches(0.35),
         "順番が崩れると、いきなり仕様の話になって経営層の興味が逃げる。",
         size=DS2.SIZE_BODY, color=DS2.WHITE, bold=True, font=DS2.FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)
footer(s)
set_notes(s, """コイケ営業マン全員に共有していただきたい 3 ステップです。
特に大事なのは順番です。仕様や機能から始めると経営層は途中で興味を失います。
まず数字 ROI、次に「実機ある、来てください」、最後に「ご予算ベースで提案します」。
この順番を社内研修にも組み込んでいただけると、TWF 終了後の長期的な営業力向上に
つながります。""")


# ============ Slide 23: アクションプラン (timeline) ============
s = new_slide(prs)
title_header(s, 23, TOTAL, "アクションプラン — 5/21 → 6/13 → 6/14 以降", kicker="ACTION PLAN")
actions = [
    ("5/21",       "本日",            "事業部長会議で討議、重点訴求 5 社決定、各営業所の重点お客様リスト作成。"),
    ("5/22-5/31",  "事前送付期間",     "重点お客様へ portal URL 送付 + 来場誘導 + 公式特設ページから事前来場登録。"),
    ("6/1-6/11",   "最終調整",         "各営業所の同行来場予定確認、当日案内ルートの事前打合せ。"),
    ("6/12-6/13",  "TWF2026 当日",     "コイケ営業マンが客先と同行来場、実機体験、商談化 / アポ獲得。"),
    ("6/14-",      "事後フォロー",      "公式パネル PDF を客先にお渡し、技術担当との詳細仕様検討へ。"),
]
y0 = Inches(2.0); row_h = Inches(0.85); gap = Inches(0.1)
for i, (date, label, body) in enumerate(actions):
    top = y0 + i * (row_h + gap)
    # left date pill (navy)
    add_rect(s, Inches(0.6), top, Inches(2.0), row_h, DS2.NAVY_DARK)
    add_text(s, Inches(0.6), top + Inches(0.12), Inches(2.0), Inches(0.35),
             date, size=DS2.SIZE_BODY + 1, color=DS2.ORANGE, bold=True, align=PP_ALIGN.CENTER, font=DS2.FONT_HEAD)
    add_text(s, Inches(0.6), top + Inches(0.48), Inches(2.0), Inches(0.3),
             label, size=DS2.SIZE_TINY, color=DS2.WHITE, align=PP_ALIGN.CENTER)
    # right body
    add_rect(s, Inches(2.75), top, DS2.SLIDE_W - Inches(3.35), row_h, DS2.WHITE, line_color=DS2.GRAY_LINE)
    add_text(s, Inches(2.95), top + Inches(0.2), DS2.SLIDE_W - Inches(3.65), row_h - Inches(0.4),
             body, size=DS2.SIZE_BODY, color=DS2.NAVY_DARK, line_spacing=1.4, anchor=MSO_ANCHOR.MIDDLE)
footer(s)
set_notes(s, """本日から TWF2026 当日までの動き方を 5 段階で整理しました。
最重要は 5/22 〜 5/31 の重点お客様への portal URL 送付です。
ここで「事前来場登録」までやっていただけると、当日の来場率が大きく変わります。
各営業所責任者にこの動き方を共有いただき、進捗を週次で確認する仕組みをお願いいたします。""")


# ============ Slide 24: クロージング (Dark sandwich の下面) ============
s = new_slide(prs)
add_rect(s, 0, 0, DS2.SLIDE_W, DS2.SLIDE_H, DS2.NAVY_DARK)
add_rect(s, 0, 0, Inches(0.6), DS2.SLIDE_H, DS2.ORANGE)
# Kicker
add_text(s, Inches(1.2), Inches(1.0), DS2.SLIDE_W - Inches(2.4), Inches(0.5),
         "CLOSING", size=DS2.SIZE_CAPTION, color=DS2.ORANGE, bold=True, font=DS2.FONT_HEAD)
# Big closing message
add_text(s, Inches(1.2), Inches(1.7), DS2.SLIDE_W - Inches(2.4), Inches(2.5),
         "TWF2026 は「展示会」ではなく\n「省人化案件創出の商談装置」",
         size=DS2.SIZE_DISPLAY - 14, color=DS2.WHITE, bold=True, font=DS2.FONT_HEAD, line_spacing=1.2)
# Divider
add_rect(s, Inches(1.2), Inches(4.6), Inches(1.0), Pt(3), DS2.ORANGE)
# Sub
add_text(s, Inches(1.2), Inches(4.85), DS2.SLIDE_W - Inches(2.4), Inches(0.6),
         "コイケ酸商 × マツモト産業 = 協業強化で「人手不足対策」の本命提案を業界に届ける",
         size=DS2.SIZE_SUBHEAD - 2, color=DS2.WHITE, font=DS2.FONT_HEAD)
# Contact section
add_text(s, Inches(1.2), Inches(5.85), DS2.SLIDE_W - Inches(2.4), Inches(0.35),
         "CONTACT", size=DS2.SIZE_CAPTION, color=DS2.ORANGE, bold=True, font=DS2.FONT_HEAD)
add_text(s, Inches(1.2), Inches(6.2), DS2.SLIDE_W - Inches(2.4), Inches(0.4),
         "マツモト産業株式会社 京葉営業所  柏原",
         size=DS2.SIZE_SUBHEAD - 2, color=DS2.WHITE, bold=True, font=DS2.FONT_HEAD)
add_text(s, Inches(1.2), Inches(6.65), DS2.SLIDE_W - Inches(2.4), Inches(0.35),
         "TEL  047-358-1121", size=DS2.SIZE_BODY, color=DS2.WHITE)
add_text(s, Inches(1.2), Inches(7.0), DS2.SLIDE_W - Inches(2.4), Inches(0.35),
         "PORTAL  twf2026-portal.pages.dev", size=DS2.SIZE_BODY, color=DS2.ORANGE, bold=True)
set_notes(s, """本日のご提案をまとめます。TWF2026 は単なる展示会ではなく、コイケ酸商様の
商談を作る装置として活用いただけます。6/12-13 までの 3 週間、コイケ営業マンと一緒に
重点お客様への portal 送付・来場誘導を進めさせていただきたいです。
本日はご清聴ありがとうございました。質疑応答に移らせていただきます。""")


# Save
prs.save(str(OUT))
print(f"saved: {OUT}")
print(f"slides: {len(prs.slides)}")
print(f"size: {OUT.stat().st_size:,} bytes")
