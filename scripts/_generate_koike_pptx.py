"""コイケ酸商 (5/21) 向け TWF2026 提案 pptx (24 枚) 生成.

設計: tmp/pptx_best_practices_research.md
デザインシステム: scripts/_pptx_design_system.py
素材: tmp/codex_official_panel_review.txt + tmp/koike_slide_data.json + data/topics.json

各スライドにスピーカーノート (200-400 字) 設定。
動画 4 箇所 (Slide 07 / 15 / 16 / 17) はプレースホルダ枠、後日 add_movie で差し込み。
"""
import sys, io
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent))
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from _pptx_design_system import (
    DS, add_rect, add_text, add_bullets, slide_title_bar, slide_footer,
    video_placeholder, set_speaker_notes,
    add_layout_title_slide, add_layout_closing_slide, add_layout_section_divider,
)

ROOT = Path("D:/repos/twf2026-portal")
OUT = ROOT / "tmp/コイケ酸商_TWF2026提案_v1.pptx"
TOTAL = 24


def add_content_slide(prs, page_no, title, subtitle, content_fn, notes):
    """Layout 3: 通常コンテンツスライド (title bar + フッター + 中身は content_fn)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_bar(slide, title, subtitle)
    content_fn(slide)
    slide_footer(slide, page_no, TOTAL)
    set_speaker_notes(slide, notes)
    return slide


def add_maker_slide(prs, page_no, maker_name, tagline, message, points, materials, portal_url, notes, video_label=None):
    """Layout 4: メーカー個別スライド (Slide 06-17)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_bar(slide, maker_name, tagline)

    # 左: 経営層メッセージ
    add_rect(slide, Inches(0.8), Inches(1.4), Inches(4.5), Inches(2.6), DS.ORANGE_TINT, line_color=DS.ORANGE)
    add_text(slide, Inches(0.95), Inches(1.55), Inches(4.2), Inches(0.4),
             "経営層メッセージ", size=11, color=DS.ORANGE_DARK, bold=True)
    add_text(slide, Inches(0.95), Inches(1.95), Inches(4.2), Inches(2.0),
             message, size=DS.SIZE_BODY, color=DS.NAVY, bold=True, line_spacing=1.45)

    # 右: TWF2026 みどころ (bullets)
    add_text(slide, Inches(5.5), Inches(1.4), Inches(7.2), Inches(0.4),
             "TWF2026 みどころ", size=12, color=DS.ORANGE_DARK, bold=True)
    add_bullets(slide, Inches(5.5), Inches(1.85), Inches(7.2), Inches(3.0),
                points, size=DS.SIZE_BODY - 1)

    # 動画プレースホルダ (該当時)
    if video_label:
        video_placeholder(slide, Inches(0.8), Inches(4.25), Inches(6.5), Inches(2.55),
                          video_label)

    # 下部: 配布資料 + portal URL
    add_rect(slide, Inches(0.8), DS.SLIDE_H - Inches(1.45), DS.SLIDE_W - Inches(1.6), Inches(0.85), DS.GRAY_BG)
    add_text(slide, Inches(1.0), DS.SLIDE_H - Inches(1.4), DS.SLIDE_W - Inches(2.0), Inches(0.35),
             f"📄 配布資料: {materials}", size=10, color=DS.GRAY_DARK)
    add_text(slide, Inches(1.0), DS.SLIDE_H - Inches(1.05), DS.SLIDE_W - Inches(2.0), Inches(0.35),
             f"🌐 詳細ページ: {portal_url}", size=10, color=DS.ORANGE_DARK, bold=True)

    slide_footer(slide, page_no, TOTAL)
    set_speaker_notes(slide, notes)
    return slide


# ====================================================================
# Build
# ====================================================================
prs = Presentation()
prs.slide_width = DS.SLIDE_W
prs.slide_height = DS.SLIDE_H


# ============ Slide 01: タイトル ============
slide = add_layout_title_slide(
    prs,
    eyebrow="TWF2026 みどころポータル",
    title_main="TWF2026 生産性向上ソリューション\nコーナーのご案内",
    subtitle="「人手不足対策」の本命企画 — 11 社一括提案",
    recipient="コイケ酸商株式会社 御中",
    presenter="マツモト産業株式会社 京葉営業所  柏原",
    date="2026 年 5 月 21 日 (水)",
)
set_speaker_notes(slide, """本日はお時間をいただきありがとうございます。
マツモト産業京葉営業所の柏原から、6 月 12-13 日に幕張メッセで開催する TWF2026
「生産性向上ソリューションコーナー」のご案内をさせていただきます。
本日のテーマは、コイケ酸商様の経営層・営業マンが、お客様の人手不足対策をご提案できる
「商談装置」として TWF2026 をご活用いただく、というご提案です。""")


# ============ Slide 02: 経営層向け結論 ============
def content_02(slide):
    add_text(slide, Inches(0.8), Inches(1.4), DS.SLIDE_W - Inches(1.6), Inches(0.6),
             "TWF2026 生産性向上ソリューションコーナーは、6 つの工程を 11 社で一括提案できる商談装置です。",
             size=16, color=DS.NAVY, bold=True)
    categories = [
        ("🤖 搬送", "ロボットバンク / シンテック"),
        ("🔥 溶接", "ファナック / ダイヘン / フロニウス / ノビテック / オートスイング"),
        ("🎨 塗装", "メサック"),
        ("🛡 安全", "小森安全機研究所"),
        ("📐 教示", "ゼネテック"),
        ("✨ クリーニング", "オプティレーザーソリューションズ"),
    ]
    positions = [
        (Inches(0.8), Inches(2.2)), (Inches(5.0), Inches(2.2)), (Inches(9.2), Inches(2.2)),
        (Inches(0.8), Inches(3.95)), (Inches(5.0), Inches(3.95)), (Inches(9.2), Inches(3.95)),
    ]
    card_w = Inches(4.0); card_h = Inches(1.55)
    for (label, makers), (left, top) in zip(categories, positions):
        add_rect(slide, left, top, card_w, card_h, DS.ORANGE_TINT, line_color=DS.ORANGE)
        add_text(slide, left + Inches(0.2), top + Inches(0.15), card_w - Inches(0.4), Inches(0.4),
                 label, size=DS.SIZE_HEADING, color=DS.ORANGE_DARK, bold=True)
        add_text(slide, left + Inches(0.2), top + Inches(0.7), card_w - Inches(0.4), Inches(0.8),
                 makers, size=12, color=DS.GRAY_DARK)
    add_text(slide, Inches(0.8), Inches(5.85), DS.SLIDE_W - Inches(1.6), Inches(0.4),
             "ご来場いただくことで、6 つの工程それぞれの省人化・高付加価値業務への再配置を",
             size=14, color=DS.GRAY_DARK)
    add_text(slide, Inches(0.8), Inches(6.25), DS.SLIDE_W - Inches(1.6), Inches(0.4),
             "目の前で確認いただけます。事業部長会議後、ご来場予約のご検討をお願いいたします。",
             size=14, color=DS.GRAY_DARK)

add_content_slide(
    prs, 2,
    title="経営層向け結論",
    subtitle="TWF2026 = 人手不足対策の本命企画、商談装置として活用",
    content_fn=content_02,
    notes="""TWF2026 は単なる展示会ではなく、「人手不足対策」を切り口に商談を作る場です。
搬送・溶接・塗装・安全・教示・クリーニングの 6 つの工程に、11 社が一括で提案を持って
集まります。コイケ酸商様のお客様が抱える「人がいない」「品質を保てない」「教育に時間が
かかる」という経営層の悩みに、TWF2026 の出展メーカーから具体的な解決策をご紹介できます。
本日はこの 11 社を 1 枚 1 社のペースでご紹介し、最後にコイケ営業マンが活用できる
アクションプランまでお話しいたします。""",
)


# ============ Slide 03: ROI 一覧 ★最重要 ============
def content_03(slide):
    roi_cards = [
        ("200%", "搬送量向上", "食品工場での AMR 導入事例", "ロボットバンク StarLift"),
        ("90%", "ティーチング時間削減", "VCOLP 採用、22 メーカー対応", "ゼネテック OLP"),
        ("47-64%", "塗料使用量削減", "自動車部品モリブデン / 光輝塗装", "メサック G05/G07"),
        ("1/4", "作業者負担を 1/4 に", "元古鉄工事例 (4 人/週 → 1 人/日)", "オプティレーザー"),
    ]
    card_w = Inches(2.85); card_h = Inches(3.7)
    for i, (number, label, evidence, source) in enumerate(roi_cards):
        left = Inches(0.8 + i * 3.0)
        top = Inches(1.5)
        add_rect(slide, left, top, card_w, card_h, DS.WHITE, line_color=DS.ORANGE)
        # 数字 (大字オレンジ、注意誘導)
        add_text(slide, left + Inches(0.1), top + Inches(0.3), card_w - Inches(0.2), Inches(1.2),
                 number, size=DS.SIZE_ROI_NUM, color=DS.ORANGE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, left + Inches(0.1), top + Inches(1.6), card_w - Inches(0.2), Inches(0.4),
                 label, size=14, color=DS.NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, left + Inches(0.2), top + Inches(2.15), card_w - Inches(0.4), Inches(1.2),
                 evidence, size=11, color=DS.GRAY_DARK, align=PP_ALIGN.CENTER, line_spacing=1.4)
        add_text(slide, left + Inches(0.2), top + Inches(3.15), card_w - Inches(0.4), Inches(0.5),
                 source, size=10, color=DS.ORANGE_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), Inches(5.55), DS.SLIDE_W - Inches(1.6), Inches(0.5),
             "※ いずれもメーカー公表値・採用事例ベース。コイケ営業マンが客先で即提案可能な数字です。",
             size=11, color=DS.GRAY_DARK)

add_content_slide(
    prs, 3,
    title="ROI 一覧 — 経営層が一番好きな数字",
    subtitle="11 社の中から、経営層向けに効く 4 つの数字を厳選",
    content_fn=content_03,
    notes="""このスライドが本日一番大事です。
搬送 200% 向上、教示 90% 削減、塗料 47-64% 削減、作業者負担 4 分の 1。
4 つとも公表値ベースの数字なので、客先で「これメーカーの公表値です」と即ご提案いただけます。
特に塗料 47-64% 削減 (メサック) と元古鉄工事例 4 倍効率 (オプティ) は経営層に強く刺さります。
コイケ営業マンが「導入したらいくら下がりますか」と聞かれた時、この数字でツカんで、
そこから個別商談に持ち込む流れが理想です。""",
)


# ============ Slide 04: コイケ取扱商品と TWF メーカーの重なり ============
def content_04(slide):
    add_text(slide, Inches(0.8), Inches(1.4), DS.SLIDE_W - Inches(1.6), Inches(0.5),
             "コイケ酸商様の主力商材と、TWF2026 生産性向上 11 社の関係を整理しました。",
             size=14, color=DS.GRAY_DARK)
    overlap = [
        ("溶接材料 / ワイヤ",     "→ フロニウス Fortis、ダイヘン VC8、ファナック CRX、ノビテック Cavitar"),
        ("ガス・ガス機器",       "→ メサック G05/G07/G08 自動ガン、塗装ブース (空気消費 35-400NL/min)"),
        ("保護具・安全用品",     "→ オートスイング WG3+ ヘルメット、小森 SRD 3D レーダー"),
        ("研磨・砥石",           "→ ファナック協働ロボパッケージ (ハンドチェンジャーでグラインダー対応)"),
        ("搬送機器",             "→ ロボットバンク StarLift / Star-7、シンテック T-Arm / Rail Station"),
        ("自動化・ティーチング", "→ ゼネテック VCOLP (22 メーカー対応)、オプティレーザー (下地処理)"),
    ]
    add_bullets(slide, Inches(0.8), Inches(2.05), DS.SLIDE_W - Inches(1.6), Inches(4.5),
                [f"● {k}\n     {v}" for k, v in overlap], size=DS.SIZE_BODY)
    add_text(slide, Inches(0.8), Inches(6.45), DS.SLIDE_W - Inches(1.6), Inches(0.45),
             "コイケのお客様の現場課題と TWF メーカーの製品が直結します。",
             size=12, color=DS.ORANGE_DARK, bold=True)

add_content_slide(
    prs, 4,
    title="コイケ取扱商品と TWF メーカーの重なり",
    subtitle="コイケのお客様 = TWF 出展メーカー製品で生産性向上できる",
    content_fn=content_04,
    notes="""コイケ酸商様の主力 6 商材と、TWF 出展 11 社の関係を整理しました。
特に溶接材料・ガス・保護具という従来のコイケ商材と、自動化・搬送・教示のような
最近の課題対応商材の橋渡しが TWF メーカーで一気に揃います。
「コイケのお客様で人手不足の現場があれば、TWF のこの 11 社で答えが出る」
という売り方ができます。""",
)


# ============ Slide 05: コイケ営業マンが活用できる場面 ============
def content_05(slide):
    scenes = [
        ("① 客先工場の課題ヒアリング時",
         "「省人化・人手不足対策」「品質安定化」「教示時間短縮」など経営層課題を聞いた瞬間に、TWF 出展 11 社の中から最適な 1-2 社を即座にご提案。"),
        ("② portal URL を送付して事前確認",
         "https://twf2026-portal.pages.dev/ をメール or LINE で送付。事業部長・担当者が事前に各社のみどころを確認した状態でご来場いただける。"),
        ("③ TWF2026 当日のご案内",
         "コイケ営業マンが客先と一緒にご来場、目の前で実機体験。塗料 47-64% 削減や搬送量 200% 向上を実演で確認 → 商談化。"),
        ("④ 当日不在でも事後フォロー",
         "公式パネル PDF を客先にお渡し。後日改めて商談時間を設定、portal の Q1-Q5 解説を読みながら導入検討を深掘り。"),
    ]
    for i, (head, body) in enumerate(scenes):
        top = Inches(1.4 + i * 1.35)
        add_rect(slide, Inches(0.8), top, DS.SLIDE_W - Inches(1.6), Inches(1.2), DS.ORANGE_TINT, line_color=DS.ORANGE)
        add_text(slide, Inches(1.0), top + Inches(0.15), DS.SLIDE_W - Inches(1.8), Inches(0.4),
                 head, size=15, color=DS.ORANGE_DARK, bold=True)
        add_text(slide, Inches(1.0), top + Inches(0.6), DS.SLIDE_W - Inches(1.8), Inches(0.55),
                 body, size=11, color=DS.GRAY_DARK)

add_content_slide(
    prs, 5,
    title="コイケ営業マンが活用できる場面",
    subtitle="客先工場でその場で提案 → 来場誘導 → 商談化",
    content_fn=content_05,
    notes="""コイケ営業マンが日常業務の中で portal をどう使えるか、4 つの場面でご紹介します。
特に大事なのは ① のヒアリング時の即提案と、② の portal 事前送付です。
portal は永続資産なので、TWF 終了後も「ヒアリングしたら portal を送る」という
習慣をつけていただけると、コイケ営業マン全員の引き出しが増えます。""",
)


# ============ Slide 06-17: メーカー 11 社個別 (うち 4 枚動画) ============

add_maker_slide(
    prs, 6, "ファナック㈱",
    "協働ロボのマルチタスク化 — 3kg 可搬 + 協働ロボパッケージ + TIG フィラー",
    "省スペース工場で 1 台を溶接・研磨・TIG に使い回すマルチタスク化。\n安全柵不要前提で固定設備化せず、小型協働ロボで小さく始められる。",
    [
        "🆕 3kg 可搬: 軽量 11kg、マグネット式溶接テーブル着脱、高電圧タッチセンサー",
        "🤖 協働ロボパッケージ: 安全柵不要、狭い工場対応",
        "🔁 ワンタッチハンドチェンジャー: CO2 トーチ ↔ グラインダー、半自動用 / TIG-研磨 / ハンドレーザ",
        "💪 力覚研磨: マツモト機械フローティング + 内蔵力覚センサ",
        "🔥 TIG フィラー仕様: 簡単教示で高品質・高能率パルス TIG",
        "🎯 連携実績: 大阪・神奈川 WF で ATC / HW1000 等を実演、TWF2026 でも継続",
    ],
    "TWF2026 公式パネル (2 ページ統合) + 3kg 可搬仕様 PDF",
    "https://twf2026-portal.pages.dev/m/fanuc/",
    notes="""ファナックは協働ロボのマルチタスク化を 4 つの切り口で提案します。
3kg 可搬は工場の好きな場所にマグネットでくっつけて使える「持ち運べる協働ロボ」、
協働ロボパッケージは安全柵なしでハンドチェンジャー付き、TIG フィラー仕様は
熟練 TIG 作業を自動化、そして過去 WF で見せてきた周辺メーカー連携も TWF2026 で
継続展示します。「固定設備でなく小さく始められる協働ロボ」という訴求が経営層に
刺さります。""",
)

add_maker_slide(
    prs, 7, "㈱ダイヘン",
    "TIG/MAG 兼用仕様 + AiTran 連携自動化 — 溶接品質の安定化",
    "1 台架台に CO2・TIG 溶接機を搭載、段取り替え簡単。\nダイヘン独自制御で高軌跡精度、難易度の高い TIG フィラーでも安定高品質。",
    [
        "🔁 TIG/MAG 兼用仕様: 1 台で 2 方式、段取り替えで使い分け",
        "🎯 高軌跡精度: TIG フィラー溶接も安定高品質",
        "🤖 AiTran 連携: 搬送→位置補正→溶接の一気通貫自動化",
        "🚀 マツモト産業ブース連動: 自動化推進コーナーで実機デモ",
        "💡 Slide 18 で「実演セミナーで同じデモが見られます」と連動訴求",
    ],
    "TWF2026 公式パネル + 実演セミナー (Slide 18 で詳細)",
    "https://twf2026-portal.pages.dev/m/daihen/",
    video_label="ダイヘン VC8 × AiTran500 連携デモ (2:29) — YouTube -ydKdIio5es",
    notes="""ダイヘンの目玉は VC8 + AiTran 連携デモです。動画は神奈川 WF で撮影した実機デモを
そのまま埋め込みます。AiTran が部材を運んで、VC8 が位置補正して溶接する、という
一連の自動化が 2 分半でわかります。重要なのは、この同じデモが TWF2026 の実演セミナー
(Slide 18) で生実演として開催される点です。「動画見た → 生で見たい」という導線で
来場動機を作ります。""",
)

add_maker_slide(
    prs, 8, "フロニウスジャパン㈱",
    "Fortis シリーズ 270〜500A — Wizard 機能で若手でも条件設定",
    "MIG/MAG・TIG・手棒・ガウジングまで 1 シリーズで対応、\nWizard 機能で経験浅でも溶接条件設定が可能、人材育成にも有効。",
    [
        "🆕 Fortis シリーズ: 270〜500A、空冷/水冷、送給装置一体型/別置き型",
        "🧙 Wizard 機能: 経験浅でも溶接条件設定が可能",
        "🔥 幅広い工法: MIG/MAG 直流・パルス、TIG、手棒、ガウジング",
        "🤖 自動化連携: ファナック CRX + TPS500i + CMT の協働ロボ連携",
        "🎁 TWF2026 ご注文特典: 自動遮光面 Vizor 4000 Plus プレゼント",
    ],
    "TWF2026 公式パネル + Manual 溶接機チラシ",
    "https://twf2026-portal.pages.dev/m/furoniusujapan/",
    notes="""フロニウスの強みは Wizard 機能です。若手作業者でも条件設定が間違いなく
できるので、人材育成と品質安定の両立に効きます。経験豊富なベテランが減る中で、
「機械側で知見を担保する」という発想は経営層に強く刺さります。
注文特典 (Vizor 4000 Plus) もありますが、これは現場担当者向けの後押し情報として
ご紹介ください。経営層プレゼンでは Wizard が主訴求です。""",
)

add_maker_slide(
    prs, 9, "ロボットバンク㈱",
    "StarLift + Star-7 — 搬送と清掃の無人化、5 業界事例",
    "AMR の 5 系統ラインナップで「人が歩く仕事」を置き換え。\n食品工場で搬送量 200% 向上、修理工場で操作教育 30 分で即戦力化。",
    [
        "🤖 StarLift 150 / 300 / 600: 積載 150〜600kg をカバー",
        "🌐 全シリーズ 5 系統: StarShip / StarMax / StarLight (低床 31cm 棚下くぐり) / RisuBot",
        "🧹 Star-7 業務用清掃ロボット: 拭き取り・掃除・吸塵・磨き (公式パネル新情報)",
        "🎯 動作仕様: 段差 20mm、登坂 8°、最小通過幅 60cm、稼働 10h、±3cm 精度",
        "👀 採用事例: 食品 200% / 修理 30 分教育 / 部品製造 導入 2 日 / 自動車部品 24 時間体制",
    ],
    "TWF2026 公式パネル + 導入事例集 + 製品ハイライト + 搬送ロボットカタログ",
    "https://twf2026-portal.pages.dev/m/robottobanku/",
    notes="""ロボットバンクは AMR 搬送ロボの 5 系統と、新発表の Star-7 清掃ロボット
の 2 本柱です。「人が歩く仕事を置き換える」という単一コンセプトで覚えていただける
のがポイントです。
食品工場での 200% 搬送量向上の事例は、コイケ酸商様のお客様の中にも食品・部品製造
の現場があれば、直接適用しやすい数字です。導入 2 日で運用開始というスピード感も
経営層の意思決定に効きます。""",
)

add_maker_slide(
    prs, 10, "㈱メサック",
    "ロボットつかみ方式塗装ブース — 1㎡省スペース + 塗料 47〜64% 削減",
    "ガンを固定しロボットがワークを持つ逆転発想で、塗装ブース自体が約 1㎡。\nG05 自動ガンで塗料使用量 47〜64% 削減事例 (自動車部品)。",
    [
        "📐 設置面積: 約 1 ㎡で省スペース設置",
        "💨 排気風量: 30 ㎥/min、ガン下向きで塗料飛散範囲限定",
        "🎯 塗料供給経路: ポンプ〜ガン間ホース 約 1m で短縮化",
        "🆕 G05/G07/G08 自動ガン: 塗料使用量 47-64% 節約事例、G08 ダイヤフラム構造 2 液塗料対応",
        "🧹 清掃性: ガン下向き固定、ロールフィルター切断式メンテで容易化",
        "🔧 工具レス分解: G08 で塗料ブロックのみ取り外し可、メンテ性大幅向上",
    ],
    "TWF2026 公式パネル + G05/G07/G08 シリーズ PDF + コンパクト塗装ブースちらし",
    "https://twf2026-portal.pages.dev/m/mesakku/",
    notes="""メサックは本日の ROI スライド (Slide 3) で出した「塗料 47-64% 削減」の出元です。
塗料は塗装現場のコストの大半を占めるので、47-64% 削減はそのまま経費削減として
試算しやすい数字です。さらに塗装ブース自体が 1㎡で済むという省スペース化と、
排気風量 30㎥/min で空調コストも下がるので、トリプルでコスト効果があります。
ロボットがワークを持って動く、という発想自体が独自で、来場時に実機で見ていただく
と納得感が違います。""",
)

add_maker_slide(
    prs, 11, "㈱ゼネテック",
    "Visual Components Robotics OLP — ティーチング 90% 削減、22 メーカー対応",
    "ワーク 3D-CAD クリックで 6 軸ロボット動作軌跡を作成。\n22 メーカーのロボットプログラム出力対応、複数メーカー混在ラインでも標準化可能。",
    [
        "🎯 教示時間 90% 削減: VCOLP 採用で 1/10",
        "🤖 22 メーカー対応: ロボットプログラム出力対象",
        "📐 CAD クリック教示: ワーク 3D-CAD クリックで 6 軸ロボット動作軌跡",
        "🆕 VCOLP 5.0: 2026 年 3 月 18 日提供開始",
        "👀 5 用途別画面: アーク/スポット/切断/研磨/塗装",
    ],
    "TWF2026 公式パネル + VCOLP パンフレット",
    "https://twf2026-portal.pages.dev/m/zenetekku/",
    notes="""ゼネテックの VCOLP は教示時間 90% 削減という強い数字です。
ロボット教示は熟練者が時間をかけて行う作業で、人件費と機会損失の両方の負担があります。
それが 1/10 になるというのは、ロボット導入の意思決定を後押しする数字です。
さらに 22 メーカー対応なので、コイケ酸商様のお客様がどのメーカーのロボットを
持っていても提案できる汎用性があります。""",
)

add_maker_slide(
    prs, 12, "㈱小森安全機研究所",
    "SRD 3D レーダー安全システム — 世界初 SIL2/PLd 規格準拠",
    "光・粉塵・煙・水・雨に強い 3D 安全レーダー。\n安全投資 = 経営価値 (労災ゼロ、設備停止リスク低減、ライン稼働率向上)。",
    [
        "🛡️ 世界初: SIL2/PLd 規格準拠 3D 安全レーダー SRD シリーズ",
        "📡 SRD 仕様: 60GHz FMCW、応答 100ms 以下、最大 6 センサ、検知距離 5m/9m/4m",
        "🌧️ 耐環境性: 光・粉塵・煙・水・雨に強く、降雨量 45mm/h 対応",
        "🎯 動的検知ゾーン: 検知ゾーン/警告ゾーンを動的設定、最大 32 種類",
        "👀 光学式が苦手な外乱環境 (溶接・塗装・搬送・屋外設備) で相談可能",
    ],
    "TWF2026 公式パネル + SRD + AIカメラ KAG 製品案内",
    "https://twf2026-portal.pages.dev/m/komori-anzen-ki-kenkyuusho/",
    notes="""小森安全機の SRD は世界初の SIL2/PLd 規格準拠 3D レーダーです。
光学センサーが苦手な現場 (粉塵・煙・水・雨) でも安全検知できるので、
溶接・塗装・搬送・屋外設備など現場の安全投資を一段上にできます。
労災ゼロ + 設備停止リスク低減 = 経営的にはライン稼働率の向上、というロジックで
ご提案ください。経営層は安全を「コスト」と捉えがちですが、「稼働率」と
言い換えると意思決定が早くなります。""",
)

add_maker_slide(
    prs, 13, "シンテック㈱",
    "3arm / T-Arm / Rail Station — 作業負荷 + 労災リスク低減",
    "重量物・工具保持・搬送補助の作業負荷を下げ、腰痛・労災リスクも低減。\nトヨタ 6000+/日野 3000+/ダイハツ 500+ セットの自動車メーカー採用実績。",
    [
        "💪 3arm: 締付・組立・バリ取り・持上操作をアシスト、最大荷重 35kg",
        "🏗️ T-Arm: 耐荷重 40〜650kg、オートバランス標準装備",
        "📈 強度向上: 引張・圧縮強度 1.8 倍に再設計",
        "🏭 採用実績: トヨタ 6000+ / 日野 3000+ / ダイハツ 500+ セット",
        "🛤️ Rail Station: 落下防止に優れた運搬搭載補助レール",
    ],
    "TWF2026 公式パネル + 3arm カタログ + 製品プレゼン",
    "https://twf2026-portal.pages.dev/m/shintech/",
    notes="""シンテックは作業負荷低減の老舗です。トヨタ 6000、日野 3000、ダイハツ 500 セット
の採用実績は信頼性の証です。
最近の課題である「ベテラン作業者の離職リスク」「腰痛による労災コスト」
「採用しても定着しない」を、3arm・T-Arm・Rail Station の組み合わせで解決できます。
省人化ではなく「働きやすさで定着率向上」というニュアンスでご提案ください。""",
)

add_maker_slide(
    prs, 14, "㈱ノビテック",
    "Cavitar Welding Camera + Weld-Eye — 溶接不良リアルタイム可視化",
    "アーク光・ヒュームをカットし溶融池を可視化、\nWeld-Eye で溶融池・面積・溶融速度を AI リアルタイム判定。",
    [
        "👀 溶融池可視化: アーク光・ヒュームをカットし、溶融池と周辺状態を詳細観察",
        "🤖 AI 解析: Weld-Eye で溶融池・面積・溶融速度を AI 解析、不良可否をリアルタイム判定",
        "🎯 2 機種比較: C350 / C400 のサイズ・fps・設置性",
        "📐 高解像度: 最大 1,440×1,080 px、C350 は最大 500 fps",
        "👀 高品位ワーク要求現場 (品質保証・原因究明・教育用途) 向け",
    ],
    "TWF2026 公式パネル + Cavitar カタログ",
    "https://twf2026-portal.pages.dev/m/nobitekku/",
    notes="""ノビテックの Cavitar Welding Camera は、アーク光に邪魔されずに溶接プールを
直接観察できるカメラです。Weld-Eye と組み合わせると AI が溶接の良否をリアルタイム判定
できるので、品質保証の証跡として動画記録できます。
高品位ワーク (航空・自動車・医療機器など) のお客様で「品質エビデンスを残したい」
というニーズがあれば即提案できます。""",
)

add_maker_slide(
    prs, 15, "㈱オートスイング (OTOS) — Ray-X 溶接カメラ",
    "WGC200 / WGC400 — アーク溶接そのものが鮮明に見える",
    "補助照明不要・超小型で、ロボット溶接やキャリッジ溶接にも設置しやすい。\n品質確認・技能教育・遠隔監視用途を具体化できる。",
    [
        "👀 Ray-X: アーク溶接そのものが鮮明、補助照明不要",
        "📷 2 機種: WGC200 (有線、152g) / WGC400 (Wi-Fi6 無線、274g)",
        "📐 計測オプション: アーク長・ビード幅・シーム位置",
        "🎬 サンプル動画: GMAW / FCAW / Orbital GTAW 撮影",
    ],
    "TWF2026 公式パネル + OTOSWING リーフレット",
    "https://twf2026-portal.pages.dev/m/ootosuingu-otos/",
    video_label="OTOS Ray-X 撮影サンプル動画 #1 (柏原所有 MP4)",
    notes="""OTOS の Ray-X は超小型の溶接カメラで、ロボット溶接にもキャリッジ溶接にも
設置しやすいのが特徴です。動画はこの後 Slide 16 と合わせて 2 種類お見せします。
本スライドは Ray-X の代表的な撮影サンプルです。GMAW、FCAW、Orbital GTAW など
複数の溶接プロセスで使えることを動画でご確認ください。""",
)

add_maker_slide(
    prs, 16, "㈱オートスイング (OTOS) — WG3+ ヘルメット",
    "カメラ搭載溶接ヘルメット — 品質確認 + 技能伝承",
    "作業者目線の可視化で、熟練者のノウハウを動画として記録・共有可能。\n後継者教育、品質トレーニング、遠隔指導に活用。",
    [
        "🪖 WG3+ ヘルメット: カメラ搭載溶接ヘルメットで作業者目線の可視化",
        "🎯 技能伝承: 熟練者のノウハウを動画で記録・共有",
        "👀 品質確認: 溶接プール・シーム位置を作業者視点で確認",
        "🎬 教育用途: 後継者教育、品質トレーニングに活用",
    ],
    "TWF2026 公式パネル + 関連製品リーフレット",
    "https://twf2026-portal.pages.dev/m/ootosuingu-otos/",
    video_label="OTOS WG3+ ヘルメット撮影サンプル動画 #2 (柏原所有 MP4)",
    notes="""WG3+ はカメラがついた溶接ヘルメットです。「熟練者がどう見ているか」を
そのまま動画で残せるので、後継者育成や品質基準の標準化に使えます。
コイケ酸商様のお客様の中で「熟練者が引退する前にノウハウを残したい」という
ニーズがあれば、保護具予算と教育予算の両方から正当化できます。""",
)

add_maker_slide(
    prs, 17, "オプティレーザーソリューションズ㈱",
    "ULT LASER シリーズ — レーザークリーナー、作業者負担 1/4",
    "錆・塗膜・酸化皮膜を非接触で除去する国内生産レーザークリーナー。\n元古鉄工事例で 4 人/1 週間 → 1 人/1 日に作業時間短縮。",
    [
        "🏭 元古鉄工事例: 4 人/1 週間 → 1 人/1 日、作業者負担 1/4",
        "🇯🇵 国内生産: 大阪本社のレーザークリーナー専門メーカー",
        "🔦 9 機種ラインナップ: CW / Pulse の用途別",
        "⚡ 即起動: 最短 5 秒、空冷 6h 以上・水冷 10h 以上の連続照射",
        "✨ 環境負荷低: 薬品・研磨材を使わず二次廃棄物を抑える",
    ],
    "TWF2026 公式パネル + 公式カタログ + 元古鉄工事例記事",
    "https://twf2026-portal.pages.dev/m/oputeireezaasoryuushonzu/",
    video_label="オプティレーザー 製品紹介 (YouTube ypxAtVayQxQ)",
    notes="""オプティレーザーは初 TWF 出展の注目株です。元古鉄工事例は ROI スライドでも
出した「4 人/1 週間 → 1 人/1 日」という強い数字の出元です。これは作業者負担 4 分の 1
+ 工期 1/5 + 副産物として二次廃棄物がほぼゼロという、3 重の効果があります。
動画では実際にレーザーで錆や塗膜を非接触で吹き飛ばす映像が見られます。
鉄骨ファブ、橋梁、船舶など重厚長大系のお客様に即訴求できます。""",
)


# ============ Slide 18: 実演セミナー紹介 ============
def content_18(slide):
    seminars = [
        ("3M (砥石・研磨)",        "🪜", "フルハーネスデモンストレーション (墜落防止実験 + 吊り下げ体験)"),
        ("神戸製鋼所",              "🔥", "AXELARC™ 新ワイヤ送給制御プロセス、建機向け安定アーク × 低スパッタ"),
        ("ダイヘン (本日 Slide 07)", "🤖", "VC8 + AiTran 連携デモ — Slide 07 で見せた動画と同内容の生実演!"),
        ("三菱電機",                "✨", "二次元ファイバレーザ加工機 ML3015GX-F60、高速・高精度切断"),
    ]
    for i, (name, icon, desc) in enumerate(seminars):
        top = Inches(1.4 + i * 1.3)
        add_rect(slide, Inches(0.8), top, DS.SLIDE_W - Inches(1.6), Inches(1.15), DS.ORANGE_TINT, line_color=DS.ORANGE)
        add_text(slide, Inches(1.0), top + Inches(0.15), Inches(1.0), Inches(0.85),
                 icon, size=36, color=DS.ORANGE_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(2.0), top + Inches(0.15), DS.SLIDE_W - Inches(2.8), Inches(0.4),
                 name, size=16, color=DS.NAVY, bold=True)
        add_text(slide, Inches(2.0), top + Inches(0.55), DS.SLIDE_W - Inches(2.8), Inches(0.6),
                 desc, size=12, color=DS.GRAY_DARK)
    add_text(slide, Inches(0.8), Inches(6.75), DS.SLIDE_W - Inches(1.6), Inches(0.4),
             "💡 Slide 07 のダイヘン動画 → 実演セミナーで生実演、という動線で来場動機を作れます。",
             size=11, color=DS.NAVY, bold=True)

add_content_slide(
    prs, 18,
    title="実演セミナー紹介 (4 社、参加無料)",
    subtitle="目の前で動く実機を見て・触れて・体感",
    content_fn=content_18,
    notes="""実演セミナーは目の前で動く実機を体感できる無料セミナーです。4 社開催で、
特にダイヘンの VC8 + AiTran 連携デモは、本日 Slide 07 でお見せした動画と
同じものが生実演されます。
「動画で見た → 実機で確かめたい」という動機付けで、コイケ酸商様のお客様の
来場誘導に効きます。3M のフルハーネス吊り下げ体験は意外と人気で、安全意識の
ある担当者が興味を持ちます。""",
)


# ============ Slide 19: 作業環境向上 + 初TWF出展 ============
def content_19(slide):
    add_rect(slide, Inches(0.8), Inches(1.4), Inches(6.0), Inches(5.4), DS.ORANGE_TINT, line_color=DS.ORANGE)
    add_text(slide, Inches(1.0), Inches(1.55), Inches(5.7), Inches(0.45),
             "作業環境向上ブース", size=DS.SIZE_HEADING, color=DS.ORANGE_DARK, bold=True)
    add_bullets(slide, Inches(1.0), Inches(2.1), Inches(5.7), Inches(4.6),
                [
                    "🌬 集塵・脱臭設備の最新提案",
                    "🦺 防塵マスク・遮光面・フルハーネス等の安全衛生製品",
                    "🔇 騒音対策、空調機器、健康診断ブース",
                    "💡 「働きやすい現場」が「採用しやすい現場」につながる経営層メッセージ",
                    "🏭 コイケのお客様の現場改善提案にも直結",
                ], size=DS.SIZE_BODY)
    add_rect(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(5.4), RGBColor_helper(0xEF, 0xF6, 0xFF), line_color=DS.NAVY)
    add_text(slide, Inches(7.2), Inches(1.55), Inches(5.2), Inches(0.45),
             "初TWF出展いちおしメーカー", size=DS.SIZE_HEADING, color=DS.NAVY, bold=True)
    add_bullets(slide, Inches(7.2), Inches(2.1), Inches(5.2), Inches(4.6),
                [
                    "🎯 オプティレーザーソリューションズ (Slide 17)",
                    "🎯 ゼネテック VCOLP (Slide 11)",
                    "🎯 ロボットバンク Star-7 (Slide 09)",
                    "🎯 他、TWF2026 で初めてのお披露目となるメーカー多数",
                    "💡 業界トレンドの先取り提案として活用可能",
                ], size=DS.SIZE_BODY)

# RGBColor を再エクスポート (content_19 で必要)
from pptx.dml.color import RGBColor as RGBColor_helper

add_content_slide(
    prs, 19,
    title="作業環境向上 + 初TWF出展いちおしメーカー",
    subtitle="生産性向上 11 社の周辺企画",
    content_fn=content_19,
    notes="""作業環境向上ブースは「働きやすい現場 = 採用しやすい現場」という切り口で
ご紹介ください。集塵・脱臭・空調・防塵マスクなど、現場改善に直結する提案が並びます。
初 TWF 出展のメーカー (オプティ・ゼネテック VCOLP 5.0・Star-7) は業界トレンドの
先取り情報として、コイケ営業マンが「他社より早い情報」を客先に届けられます。""",
)


# ============ Slide 20: portal 紹介 ============
def content_20(slide):
    add_rect(slide, Inches(0.8), Inches(1.4), Inches(7.5), Inches(4.0), DS.ORANGE_TINT, line_color=DS.ORANGE)
    add_text(slide, Inches(1.0), Inches(1.55), Inches(7.1), Inches(0.5),
             "https://twf2026-portal.pages.dev/", size=DS.SIZE_HEADING, color=DS.ORANGE_DARK, bold=True)
    portal_features = [
        "● 149 出展メーカーの portal、生産性向上 11 社を含む全社情報を収録",
        "● 各社の Q1-Q5 (企画概要・新製品・みどころ・特典・配布資料) を整理",
        "● 公式パネル PDF・カタログ・チラシをワンクリックで DL 可能",
        "● スマホ対応、来場前・来場後どちらも快適に閲覧",
        "● 事前来場登録ボタン → 公式特設ページ経由で 1 分登録",
        "● 5/18 時点で生産性向上 11 社の Q1-Q5 が経営層向けに完全リライト済",
    ]
    add_bullets(slide, Inches(1.0), Inches(2.1), Inches(7.1), Inches(3.2),
                portal_features, size=DS.SIZE_BODY)
    # QR コードプレースホルダ
    add_rect(slide, Inches(8.7), Inches(1.4), Inches(3.8), Inches(4.0), DS.WHITE, line_color=DS.NAVY)
    add_text(slide, Inches(8.8), Inches(1.55), Inches(3.6), Inches(0.4),
             "QR コード", size=11, color=DS.GRAY_DARK, align=PP_ALIGN.CENTER, bold=True)
    add_rect(slide, Inches(9.5), Inches(2.1), Inches(2.2), Inches(2.2), DS.NAVY)
    add_text(slide, Inches(9.5), Inches(2.5), Inches(2.2), Inches(0.5),
             "QR", size=36, color=DS.WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(9.5), Inches(3.2), Inches(2.2), Inches(0.9),
             "(配布版で QR\n挿入予定)", size=10, color=DS.WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(8.8), Inches(4.4), Inches(3.6), Inches(0.9),
             "スマホで二次元コード\n読取り即アクセス", size=11, color=DS.GRAY_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), Inches(5.7), DS.SLIDE_W - Inches(1.6), Inches(0.5),
             "コイケ営業マンが客先と一緒にスマホで開いて、その場で工程別の見どころをご案内できます。",
             size=12, color=DS.NAVY)

add_content_slide(
    prs, 20,
    title="TWF2026 みどころポータル",
    subtitle="コイケ営業マンが普段から使える資料",
    content_fn=content_20,
    notes="""portal はコイケ営業マンの「普段使い資料」です。客先でスマホを開いて、
お客様と一緒に見ながらご案内する想定で設計してあります。
特に強調したいのは、5/18 時点で生産性向上 11 社の Q1-Q5 が経営層向けに完全リライト済
である点です。営業マンがそのまま読み上げてもプレゼンになるレベルで整理しています。
配布版にはこのページに QR コードを差し込みます。""",
)


# ============ Slide 21: 公式パネル PDF + パンフレットの使い方 ============
def content_21(slide):
    rows = [
        ("公式パネル PDF (10 社分)", "ブース前説明・営業配布用の一次資料",
         "経営層向けの「掴み」、客先で 30 秒〜1 分で概要を伝える。各社 1-2 ページの簡潔な内容、PDF で印刷も可。"),
        ("既存パンフレット",         "仕様確認・商談深掘り・導入検討用",
         "各社 10-30 ページの詳細仕様。技術担当との具体的検討、見積前提条件の整理に使用。portal から DL 可能。"),
        ("3kg 可搬仕様 PDF (ファナック)", "新製品の単独訴求",
         "ファナック最新協働ロボの 1 ページ簡易資料。デモ依頼受付ありの旨を明示。"),
    ]
    top0 = Inches(1.5)
    for i, (col1, col2, col3) in enumerate(rows):
        top = top0 + Inches(i * 1.65)
        add_rect(slide, Inches(0.8), top, Inches(3.5), Inches(1.5), DS.ORANGE_TINT, line_color=DS.ORANGE)
        add_text(slide, Inches(0.95), top + Inches(0.2), Inches(3.2), Inches(0.4),
                 col1, size=14, color=DS.ORANGE_DARK, bold=True)
        add_text(slide, Inches(0.95), top + Inches(0.65), Inches(3.2), Inches(0.8),
                 col2, size=11, color=DS.NAVY, bold=True)
        add_rect(slide, Inches(4.5), top, DS.SLIDE_W - Inches(5.3), Inches(1.5), DS.WHITE, line_color=DS.GRAY_LIGHT)
        add_text(slide, Inches(4.65), top + Inches(0.2), DS.SLIDE_W - Inches(5.55), Inches(1.2),
                 col3, size=12, color=DS.GRAY_DARK, line_spacing=1.35)

add_content_slide(
    prs, 21,
    title="資料の使い方 — 公式パネル PDF と既存パンフレット",
    subtitle="場面別の使い分けで商談効率アップ",
    content_fn=content_21,
    notes="""3 種類の資料の使い分けを整理しました。
公式パネル PDF は経営層への「掴み」用、各社 1-2 ページで簡潔。
既存パンフレットは技術担当との詳細検討用。
ファナック 3kg 可搬は新製品の単独訴求用です。
コイケ営業マンが客先で「まず何を渡すか」迷ったときの判断基準として
ご活用ください。""",
)


# ============ Slide 22: 営業トーク順 ============
def content_22(slide):
    steps = [
        ("Step 1: まず数字 ROI でツカむ",
         "「塗料 47-64% 削減」「搬送量 200% 向上」「ティーチング 90% 削減」を最初に出す。経営層は数字に反応する。"),
        ("Step 2: 実機体験へ TWF 来場誘導",
         "「実機を目の前で動かしてます、6/12-13 に幕張メッセでぜひご確認を」と TWF 来場を提案。公式特設ページから事前登録。"),
        ("Step 3: 個別相談へ誘導 (商談化)",
         "「導入規模・タイミング・ご予算感をお聞かせいただければ、メーカーと一緒に最適構成をご提案します」と商談化。"),
    ]
    for i, (head, body) in enumerate(steps):
        top = Inches(1.4 + i * 1.65)
        add_rect(slide, Inches(0.8), top, DS.SLIDE_W - Inches(1.6), Inches(1.5), DS.ORANGE_TINT, line_color=DS.ORANGE)
        add_text(slide, Inches(1.0), top + Inches(0.2), DS.SLIDE_W - Inches(1.8), Inches(0.5),
                 head, size=16, color=DS.ORANGE_DARK, bold=True)
        add_text(slide, Inches(1.0), top + Inches(0.75), DS.SLIDE_W - Inches(1.8), Inches(0.7),
                 body, size=12, color=DS.GRAY_DARK, line_spacing=1.35)
    add_text(slide, Inches(0.8), Inches(6.5), DS.SLIDE_W - Inches(1.6), Inches(0.4),
             "💡 この順番が崩れると、いきなり仕様の話になって経営層の興味が逃げます。",
             size=11, color=DS.NAVY, bold=True)

add_content_slide(
    prs, 22,
    title="営業トーク順 — 数字 → 実機体験 → 個別相談",
    subtitle="3 ステップで商談化",
    content_fn=content_22,
    notes="""コイケ営業マン全員に共有していただきたい 3 ステップです。
特に大事なのは順番です。仕様や機能から始めると経営層は途中で興味を失います。
まず数字 ROI、次に「実機ある、来てください」、最後に「ご予算ベースで提案します」。
この順番を社内研修にも組み込んでいただけると、TWF 終了後の長期的な営業力
向上にもつながります。""",
)


# ============ Slide 23: コイケ来場アクションプラン ============
def content_23(slide):
    actions = [
        ("📅 5/21 (本日)",         "本資料をもとに事業部長会議で討議。重点訴求 5 社を決定、各営業所の重点お客様リストを作成。"),
        ("📅 5/22 〜 5/31",        "重点お客様への portal URL 送付 + 来場誘導 (メール / LINE / 訪問)。公式特設ページから事前来場登録の案内も同時実施。"),
        ("📅 6/1 〜 6/11",         "各営業所の同行来場予定の最終確認。当日案内ルートの事前打合せ。"),
        ("📅 6/12-13 (TWF2026)",   "コイケ営業マンが客先と同行来場、目の前で実機体験。商談化、または後日訪問のアポ獲得。"),
        ("📅 6/14 〜",             "来場後フォロー: 公式パネル PDF を客先にお渡し、技術担当との詳細仕様検討へ。"),
    ]
    for i, (date, body) in enumerate(actions):
        top = Inches(1.4 + i * 1.05)
        add_rect(slide, Inches(0.8), top, Inches(2.5), Inches(0.9), DS.ORANGE_TINT, line_color=DS.ORANGE)
        add_text(slide, Inches(0.95), top + Inches(0.15), Inches(2.2), Inches(0.6),
                 date, size=13, color=DS.ORANGE_DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, Inches(3.5), top, DS.SLIDE_W - Inches(4.3), Inches(0.9), DS.WHITE, line_color=DS.GRAY_LIGHT)
        add_text(slide, Inches(3.7), top + Inches(0.15), DS.SLIDE_W - Inches(4.7), Inches(0.6),
                 body, size=11, color=DS.GRAY_DARK, line_spacing=1.3)

add_content_slide(
    prs, 23,
    title="コイケ来場アクションプラン",
    subtitle="5/21 → 6/12-13 までの動き方",
    content_fn=content_23,
    notes="""本日から TWF2026 当日までの動き方を 5 段階で整理しました。
最重要は 5/22 〜 5/31 の重点お客様への portal URL 送付です。
ここで「事前来場登録」までやっていただけると、当日の来場率が
大きく変わります。各営業所責任者にこの動き方を共有いただき、
進捗を週次で確認する仕組みをお願いいたします。""",
)


# ============ Slide 24: クロージング ============
slide = add_layout_closing_slide(
    prs,
    headline="TWF2026 は「展示会」ではなく\n「省人化案件創出の商談装置」です",
    sub_message="コイケ酸商様 + マツモト産業 = 協業強化で「人手不足対策」の本命提案を業界に届ける",
    contact_lines=[
        "マツモト産業株式会社 京葉営業所  柏原",
        "TEL: 047-358-1121",
        "TWF2026 みどころポータル: https://twf2026-portal.pages.dev/",
    ],
)
set_speaker_notes(slide, """本日のご提案をまとめます。
TWF2026 は単なる展示会ではなく、コイケ酸商様の商談を作る装置として活用いただけます。
6/12-13 までの 3 週間、コイケ営業マンと一緒に重点お客様への portal 送付・来場誘導
を進めさせていただきたいです。
本日はご清聴ありがとうございました。質疑応答に移らせていただきます。""")


# ====================================================================
prs.save(str(OUT))
print(f"saved: {OUT}")
print(f"slides: {len(prs.slides)}")
print(f"size: {OUT.stat().st_size:,} bytes")
