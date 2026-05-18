"""コイケ酸商向け pptx v2 デザインシステム (McKinsey 風 + portal オレンジ).

公式 pptx skill SKILL.md と McKinsey ブランドガイドの両方を取り込んだ v2。

主要原則:
1. Dominance over equality: White 60% + Navy 30% + Orange 10% (アクセントは sparingly)
2. Dark/light sandwich: 表紙 + クロージング dark、コンテンツ light
3. Sharp corners (border radius 0) — McKinsey 美意識
4. Strong contrast: Dark navy text on white、white text on navy
5. Visual motif: シャープな矩形ストリップ + 大型 stat callouts (60-72pt 数字)
6. Text-only NG: 各スライドに visual element 必須
7. アクセントライン下線禁止 (AI 生成スライドの典型!)
8. Title 36-44pt bold、Section 20-24pt、Body 14-16pt
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class DS2:
    """McKinsey 風 + portal オレンジ Design System v2."""
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    # カラー (McKinsey ベース、blue → portal orange に置換)
    NAVY_DARK = RGBColor(0x05, 0x1C, 0x2C)     # McKinsey #051C2C 相当、本文・タイトル
    NAVY = RGBColor(0x1F, 0x29, 0x37)          # portal ネイビー、見出し
    ORANGE = RGBColor(0xF9, 0x73, 0x16)        # portal オレンジ、accent sparingly
    ORANGE_DARK = RGBColor(0xC2, 0x57, 0x0C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY_BG = RGBColor(0xF7, 0xF8, 0xFA)       # 薄い背景
    GRAY_LINE = RGBColor(0xD1, 0xD5, 0xDB)     # 区切り線
    GRAY_TEXT = RGBColor(0x4B, 0x55, 0x63)     # 補助テキスト

    # フォント: 公式 SKILL.md 推奨ペアの "Calibri" + Yu Gothic fallback
    # PowerPoint で日本語が自動で Yu Gothic に切替わるよう Calibri を指定
    FONT_BODY = "Calibri"
    FONT_HEAD = "Calibri"
    FONT_JP = "Yu Gothic"  # 日本語明示時用

    # サイズ (公式 SKILL.md 推奨 + McKinsey H2 44pt)
    SIZE_DISPLAY = 64       # 表紙メイン、クロージングメッセージ
    SIZE_STAT = 72          # 大型 stat callouts (公式 large stat callouts 60-72pt 上限)
    SIZE_STAT_M = 56        # 中型 stat
    SIZE_TITLE = 36         # スライドタイトル (公式 36-44pt)
    SIZE_SUBHEAD = 20       # サブタイトル
    SIZE_SECTION = 22       # セクション見出し
    SIZE_BODY = 15          # 本文 (公式 14-16pt)
    SIZE_CAPTION = 11       # キャプション
    SIZE_TINY = 9


def add_rect(slide, left, top, width, height, fill, line_color=None, line_width=None):
    """シャープ矩形 (McKinsey 美意識: border radius 0)。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text,
             size, color, bold=False, font=DS2.FONT_BODY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items,
                size=DS2.SIZE_BODY, color=DS2.NAVY_DARK, font=DS2.FONT_BODY,
                line_spacing=1.4, space_after=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return tb


def title_header(slide, page_num, total, title, kicker=None):
    """McKinsey 風タイトル: 左上 page no + 中央タイトル、アクセント下線 NO (AI 典型回避)。
    代わりに左端 0.2in オレンジ縦バー (visual motif) + 余白で表現。
    """
    # 左端オレンジ縦バー (visual motif、全スライド共通)
    add_rect(slide, Inches(0), Inches(0), Inches(0.3), DS2.SLIDE_H, DS2.ORANGE)
    # ページ番号 (左上)
    add_text(slide, Inches(0.6), Inches(0.35), Inches(2), Inches(0.3),
             f"{page_num:02d} / {total:02d}", size=DS2.SIZE_TINY, color=DS2.GRAY_TEXT, bold=True)
    # Kicker (オレンジ small text)
    if kicker:
        add_text(slide, Inches(0.6), Inches(0.7), DS2.SLIDE_W - Inches(1.2), Inches(0.35),
                 kicker, size=DS2.SIZE_CAPTION, color=DS2.ORANGE_DARK, bold=True, font=DS2.FONT_HEAD)
    # Title (大、bold、navy)
    title_top = Inches(1.0) if kicker else Inches(0.6)
    add_text(slide, Inches(0.6), title_top, DS2.SLIDE_W - Inches(1.2), Inches(1.0),
             title, size=DS2.SIZE_TITLE, color=DS2.NAVY_DARK, bold=True,
             font=DS2.FONT_HEAD, line_spacing=1.15)
    # ※ アクセント下線は引かない (AI 生成スライドの典型回避)


def footer(slide, brand_text="マツモト産業㈱ 京葉営業所 / TWF2026 みどころポータル"):
    """フッター: 左下にブランド、右下に URL。区切り線は薄いグレー。"""
    add_rect(slide, Inches(0.6), DS2.SLIDE_H - Inches(0.55), DS2.SLIDE_W - Inches(1.2), Pt(0.5), DS2.GRAY_LINE)
    add_text(slide, Inches(0.6), DS2.SLIDE_H - Inches(0.45),
             Inches(8), Inches(0.35),
             brand_text, size=DS2.SIZE_TINY, color=DS2.GRAY_TEXT)
    add_text(slide, DS2.SLIDE_W - Inches(4), DS2.SLIDE_H - Inches(0.45),
             Inches(3.4), Inches(0.35),
             "twf2026-portal.pages.dev", size=DS2.SIZE_TINY, color=DS2.ORANGE_DARK,
             bold=True, align=PP_ALIGN.RIGHT)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def stat_callout(slide, left, top, width, height, number, label, evidence, source):
    """McKinsey 風 stat block: 大数字 + 簡潔ラベル + エビデンス + 出典.
    visual motif: 上側オレンジ細帯 + 白背景。
    """
    # 上側オレンジ帯 (Visual motif)
    add_rect(slide, left, top, width, Inches(0.15), DS2.ORANGE)
    add_rect(slide, left, top + Inches(0.15), width, height - Inches(0.15), DS2.WHITE,
             line_color=DS2.GRAY_LINE)
    # 大数字 (60-72pt、navy で迫力)
    add_text(slide, left, top + Inches(0.4), width, Inches(1.5),
             number, size=DS2.SIZE_STAT, color=DS2.NAVY_DARK, bold=True,
             align=PP_ALIGN.CENTER, font=DS2.FONT_HEAD, line_spacing=1.0)
    # ラベル (orange、small)
    add_text(slide, left, top + Inches(1.95), width, Inches(0.45),
             label, size=DS2.SIZE_SECTION - 4, color=DS2.ORANGE_DARK, bold=True,
             align=PP_ALIGN.CENTER, font=DS2.FONT_HEAD)
    # エビデンス
    add_text(slide, left + Inches(0.2), top + Inches(2.5), width - Inches(0.4), Inches(1.0),
             evidence, size=DS2.SIZE_CAPTION, color=DS2.GRAY_TEXT,
             align=PP_ALIGN.CENTER, line_spacing=1.35)
    # 出典 (orange small)
    add_text(slide, left + Inches(0.2), top + height - Inches(0.5), width - Inches(0.4), Inches(0.4),
             source, size=DS2.SIZE_TINY, color=DS2.ORANGE_DARK, bold=True,
             align=PP_ALIGN.CENTER, font=DS2.FONT_HEAD)


def two_col_card(slide, left, top, width, height, title, body, accent=False):
    """2 列レイアウト用カード."""
    fill = DS2.WHITE
    border = DS2.GRAY_LINE
    add_rect(slide, left, top, width, height, fill, line_color=border)
    if accent:
        # 上側オレンジ帯
        add_rect(slide, left, top, width, Inches(0.1), DS2.ORANGE)
    add_text(slide, left + Inches(0.25), top + Inches(0.25), width - Inches(0.5), Inches(0.45),
             title, size=DS2.SIZE_SECTION - 2, color=DS2.NAVY_DARK, bold=True,
             font=DS2.FONT_HEAD)
    add_text(slide, left + Inches(0.25), top + Inches(0.8), width - Inches(0.5), height - Inches(1.1),
             body, size=DS2.SIZE_BODY - 1, color=DS2.GRAY_TEXT, line_spacing=1.4)


def video_frame(slide, left, top, width, height, label):
    """動画埋込予定枠 (16:9 推奨)。Dark + オレンジ枠で「ここに動画」が一目で分かる。"""
    add_rect(slide, left, top, width, height, DS2.NAVY_DARK, line_color=DS2.ORANGE)
    # オレンジ play 三角風
    add_text(slide, left, top + height/2 - Inches(0.45), width, Inches(0.5),
             "▶  動画埋込予定", size=24, color=DS2.ORANGE, bold=True,
             align=PP_ALIGN.CENTER, font=DS2.FONT_HEAD)
    add_text(slide, left, top + height/2 + Inches(0.1), width, Inches(0.4),
             label, size=DS2.SIZE_CAPTION, color=DS2.WHITE, align=PP_ALIGN.CENTER)
